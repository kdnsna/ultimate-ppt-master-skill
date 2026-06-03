#!/usr/bin/env python3
"""Desktop worker for the Ultimate PPT Master app.

The worker creates a local project folder, normalizes the user's source into a
first-pass outline, and emits a production draft that is strong enough for local
review. Full agentic refinement still lives in SKILL.md and scripts/, but the
desktop output should use the same design assets and visual language instead of
a throwaway smoke-test shell.
"""

from __future__ import annotations

import argparse
import contextlib
import html
import importlib.util
import io
import json
import os
import shlex
import platform
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover - covered by environment checks
    Presentation = None  # type: ignore
    RGBColor = None  # type: ignore
    MSO_SHAPE = None  # type: ignore
    PP_ALIGN = None  # type: ignore
    Inches = None  # type: ignore
    Pt = None  # type: ignore


OUTPUT_MODES = {"pptx", "web"}
STYLE_PRESETS = {"business", "consulting", "academic", "editorial", "swiss"}
SOURCE_KINDS = {"file", "text", "url", "markdown"}
TEXT_SUFFIXES = {".md", ".markdown", ".txt", ".csv", ".tsv", ".json"}
DOCUMENT_SUFFIXES = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}
DEFAULT_PROJECT_ROOT = "projects/desktop"
ENV_FILE_NAME = ".env"
USER_ENV_FILE = Path.home() / ".ppt-master" / ".env"
PROVIDER_KEY_GROUPS = {
    "openai": ("OPENAI_API_KEY",),
    "gemini": ("GEMINI_API_KEY",),
    "qwen": ("QWEN_API_KEY", "DASHSCOPE_API_KEY"),
    "deepseek": ("DEEPSEEK_API_KEY",),
    "pexels": ("PEXELS_API_KEY",),
    "pixabay": ("PIXABAY_API_KEY",),
    "elevenlabs": ("ELEVENLABS_API_KEY",),
    "minimax": ("MINIMAX_API_KEY",),
    "cosyvoice": ("COSYVOICE_API_KEY", "DASHSCOPE_API_KEY"),
}
DIRECT_LLM_KEYS = ("LLM_API_KEY", "OPENAI_API_KEY", "GEMINI_API_KEY", "QWEN_API_KEY", "DASHSCOPE_API_KEY", "DEEPSEEK_API_KEY")
MODEL_PROVIDERS = {"auto", "openai", "gemini", "qwen", "deepseek", "custom"}
IMAGE_PROVIDERS = {"auto", "openai", "gemini", "qwen", "pexels", "pixabay", "none"}
VOICE_PROVIDERS = {"edge", "elevenlabs", "minimax", "qwen", "cosyvoice", "none"}
DEFAULT_PROVIDER_CONFIG = {
    "modelProvider": "auto",
    "textModelId": "",
    "imageProvider": "auto",
    "imageModelId": "",
    "narrationEnabled": False,
    "voiceProvider": "edge",
    "voiceId": "zh-CN-XiaoxiaoNeural",
    "voiceRate": "+0%",
}


@dataclass(frozen=True)
class Step:
    key: str
    label: str
    message: str
    progress: int


STEPS = [
    Step("source", "资料处理", "读取源资料并创建本地项目。", 14),
    Step("strategy", "策略整理", "提取标题、关键句和页面骨架。", 32),
    Step("design", "设计锁定", "应用所选输出模式和风格预设。", 50),
    Step("generate", "生成预览", "生成可预览的 PPTX 或网页演示文件。", 74),
    Step("verify", "校验", "检查输出路径、预览内容和日志。", 90),
    Step("export", "导出", "保存结果并写入项目清单。", 100),
]


def sanitize_name(value: str) -> str:
    value = value.strip() or "untitled"
    safe = "".join(ch if ch.isalnum() or ch in "-_." else "-" for ch in value)
    safe = re.sub(r"-{2,}", "-", safe).strip("-_.")
    return safe[:80] or "untitled"


def repo_root_from_args(value: str | None) -> Path:
    if value:
        return Path(value).expanduser().resolve()
    return Path(__file__).resolve().parents[3]


def read_json_from_stdin() -> dict[str, Any]:
    raw = sys.stdin.read()
    if not raw.strip():
        raise ValueError("No JSON job payload received on stdin.")
    value = json.loads(raw)
    if not isinstance(value, dict):
        raise ValueError("Job payload must be a JSON object.")
    return value


def clean_optional_text(value: Any, *, max_len: int = 160) -> str:
    if value is None:
        return ""
    return str(value).strip()[:max_len]


def normalize_provider_config(raw: Any) -> dict[str, Any]:
    config = dict(DEFAULT_PROVIDER_CONFIG)
    if isinstance(raw, dict):
        model_provider = clean_optional_text(raw.get("modelProvider") or raw.get("model_provider")).lower()
        image_provider = clean_optional_text(raw.get("imageProvider") or raw.get("image_provider")).lower()
        voice_provider = clean_optional_text(raw.get("voiceProvider") or raw.get("voice_provider")).lower()
        if model_provider in MODEL_PROVIDERS:
            config["modelProvider"] = model_provider
        if image_provider in IMAGE_PROVIDERS:
            config["imageProvider"] = image_provider
        if voice_provider in VOICE_PROVIDERS:
            config["voiceProvider"] = voice_provider
        config["textModelId"] = clean_optional_text(raw.get("textModelId") or raw.get("text_model_id"))
        config["imageModelId"] = clean_optional_text(raw.get("imageModelId") or raw.get("image_model_id"))
        config["voiceId"] = clean_optional_text(raw.get("voiceId") or raw.get("voice_id")) or config["voiceId"]
        config["voiceRate"] = clean_optional_text(raw.get("voiceRate") or raw.get("voice_rate"), max_len=16) or "+0%"
        config["narrationEnabled"] = bool(raw.get("narrationEnabled", raw.get("narration_enabled", False)))
    if config["voiceProvider"] == "none":
        config["narrationEnabled"] = False
    if not config["narrationEnabled"]:
        config["voiceProvider"] = "none"
        config["voiceId"] = ""
        config["voiceRate"] = ""
    return config


def validate_job(job: dict[str, Any]) -> dict[str, Any]:
    source = job.get("source")
    if not isinstance(source, dict):
        raise ValueError("job.source is required.")

    source_kind = str(source.get("kind", "")).lower()
    source_value = str(source.get("value", "")).strip()
    if source_kind not in SOURCE_KINDS:
        raise ValueError(f"Unsupported source kind: {source_kind}")
    if not source_value:
        raise ValueError("source.value is required.")

    output_mode = str(job.get("outputMode", job.get("output_mode", ""))).lower()
    if output_mode not in OUTPUT_MODES:
        raise ValueError(f"Unsupported output mode: {output_mode}")

    style_preset = str(job.get("stylePreset", job.get("style_preset", ""))).lower()
    if style_preset not in STYLE_PRESETS:
        raise ValueError(f"Unsupported style preset: {style_preset}")

    project_dir = job.get("projectDir", job.get("project_dir"))
    if project_dir is not None and not isinstance(project_dir, str):
        raise ValueError("projectDir must be a string when provided.")

    return {
        "source": {
            "kind": source_kind,
            "value": source_value,
            "name": str(source.get("name") or "").strip(),
        },
        "outputMode": output_mode,
        "stylePreset": style_preset,
        "projectDir": project_dir,
        "providerConfig": normalize_provider_config(job.get("providerConfig", job.get("provider_config"))),
    }


def utc_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def source_suffix(source: dict[str, Any]) -> str:
    name = source.get("name") or source.get("value") or ""
    if source.get("kind") == "url":
        return ".url"
    return Path(str(name)).suffix.lower()


def recommend_settings(source: dict[str, Any]) -> dict[str, Any]:
    kind = str(source.get("kind", "")).lower()
    value = str(source.get("value", ""))
    name = str(source.get("name", ""))
    suffix = source_suffix(source)
    text = f"{name}\n{value}".lower()

    if kind == "url":
        return {
            "outputMode": "web",
            "stylePreset": "editorial",
            "pageRange": "6-10",
            "reason": "URL 内容通常更适合先做成可浏览的网页演示，便于快速预览和分享。",
        }
    if suffix in {".pptx", ".docx", ".xlsx", ".xlsm", ".pdf"}:
        return {
            "outputMode": "pptx",
            "stylePreset": "business",
            "pageRange": "8-14",
            "reason": "正式资料优先生成真实可编辑 PPTX，便于交付后继续修改。",
        }
    if "swiss" in text or "工程" in text or "product" in text or "架构" in text:
        return {
            "outputMode": "web",
            "stylePreset": "swiss",
            "pageRange": "7-12",
            "reason": "产品、工程和系统表达适合 Swiss Style 的网格、秩序和信息密度。",
        }
    if "演讲" in text or "发布" in text or "分享" in text or "keynote" in text:
        return {
            "outputMode": "web",
            "stylePreset": "editorial",
            "pageRange": "8-12",
            "reason": "演讲和发布场景更需要视觉记忆点，网页演示更有现场感。",
        }
    if "咨询" in text or "方案" in text or "strategy" in text:
        return {
            "outputMode": "pptx",
            "stylePreset": "consulting",
            "pageRange": "10-16",
            "reason": "方案材料适合结论先行的咨询风，并保留 PPTX 可编辑交付。",
        }
    if "课程" in text or "培训" in text or "学术" in text or "论文" in text:
        return {
            "outputMode": "pptx",
            "stylePreset": "academic",
            "pageRange": "12-20",
            "reason": "培训和学术内容需要章节清楚、证据稳定，PPTX 更适合二次讲授。",
        }
    return {
        "outputMode": "pptx",
        "stylePreset": "business",
        "pageRange": "8-12",
        "reason": "默认走可编辑 PPTX，覆盖多数正式汇报和团队交付场景。",
    }


def command_exists(name: str) -> bool:
    return shutil.which(name) is not None


def python_module_available(module_name: str) -> bool:
    return importlib.util.find_spec(module_name) is not None


def strip_env_value(value: str) -> str:
    value = value.strip()
    if not value:
        return ""
    quote = value[0]
    if quote in {"'", '"'}:
        end = value.find(quote, 1)
        if end > 0:
            return value[1:end].strip()
    return value.split(" #", 1)[0].strip().strip("'\"")


def resolve_env_file(repo_root: Path) -> Path | None:
    for candidate in (Path.cwd() / ENV_FILE_NAME, repo_root / ENV_FILE_NAME, USER_ENV_FILE):
        if candidate.exists():
            return candidate
    return None


def read_env_file_flags(path: Path | None) -> dict[str, str]:
    if path is None:
        return {}
    values: dict[str, str] = {}
    try:
        lines = path.read_text(encoding="utf-8").splitlines()
    except OSError:
        return values

    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith("#") or "=" not in stripped:
            continue
        key, raw_value = stripped.split("=", 1)
        key = key.strip()
        if key:
            values[key] = strip_env_value(raw_value)
    return values


def env_or_file_has(key: str, file_values: dict[str, str]) -> bool:
    return bool(os.environ.get(key) or file_values.get(key))


def env_or_file_value(key: str, file_values: dict[str, str]) -> str:
    return os.environ.get(key) or file_values.get(key, "")


def cairo_available() -> bool:
    if command_exists("pkg-config"):
        result = subprocess.run(
            ["pkg-config", "--exists", "cairo"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            check=False,
        )
        if result.returncode == 0:
            return True
    return command_exists("cairo-trace")


def inspect_environment(repo_root: Path) -> dict[str, Any]:
    venv_python = repo_root / ".venv" / "bin" / "python"
    env_file = resolve_env_file(repo_root)
    env_file_values = read_env_file_flags(env_file)
    llm_provider = env_or_file_value("LLM_PROVIDER", env_file_values) or env_or_file_value("LLM_DRIVER", env_file_values)
    llm_model = env_or_file_value("LLM_MODEL", env_file_values)
    direct_llm_configured = bool((llm_provider or llm_model) and any(env_or_file_has(key, env_file_values) for key in DIRECT_LLM_KEYS))
    return {
        "repoRoot": str(repo_root),
        "platform": platform.platform(),
        "python": {
            "executable": sys.executable,
            "version": platform.python_version(),
            "bundledVenv": venv_python.exists(),
            "pythonPptx": Presentation is not None,
        },
        "node": {
            "available": command_exists("node"),
            "npm": command_exists("npm"),
            "pnpm": command_exists("pnpm"),
        },
        "optional": {
            "cairo": cairo_available(),
            "rust": command_exists("rustc") and command_exists("cargo"),
            "edgeTts": python_module_available("edge_tts"),
        },
        "providers": {
            key: any(env_or_file_has(env_key, env_file_values) for env_key in env_keys)
            for key, env_keys in PROVIDER_KEY_GROUPS.items()
        },
        "config": {
            "envFile": str(env_file) if env_file else None,
            "imageBackend": env_or_file_value("IMAGE_BACKEND", env_file_values) or None,
            "llmProvider": llm_provider or None,
            "llmModel": llm_model or None,
            "directLlmConfigured": direct_llm_configured,
        },
    }


def model_provider_ready(env: dict[str, Any], config: dict[str, Any]) -> bool:
    provider = str(config.get("modelProvider") or "auto")
    if provider == "auto":
        return any(env["providers"].get(name) for name in ("openai", "gemini", "qwen", "deepseek"))
    if provider == "custom":
        return bool(config.get("textModelId"))
    return bool(env["providers"].get(provider))


def image_provider_ready(env: dict[str, Any], config: dict[str, Any]) -> bool:
    provider = str(config.get("imageProvider") or "auto")
    if provider == "none":
        return True
    if provider == "auto":
        return any(env["providers"].get(name) for name in ("openai", "gemini", "qwen", "pexels", "pixabay"))
    return bool(env["providers"].get(provider))


def voice_provider_ready(env: dict[str, Any], config: dict[str, Any]) -> bool:
    if not config.get("narrationEnabled"):
        return True
    provider = str(config.get("voiceProvider") or "edge")
    if provider == "edge":
        return bool(env["optional"].get("edgeTts"))
    if provider == "qwen":
        return bool(env["providers"].get("qwen"))
    return bool(env["providers"].get(provider))


def build_project_checks(env: dict[str, Any], job: dict[str, Any], source_name: str) -> list[dict[str, str]]:
    suffix = Path(source_name).suffix.lower()
    provider_config = job["providerConfig"]
    model_provider = provider_config["modelProvider"]
    image_provider = provider_config["imageProvider"]
    voice_provider = provider_config["voiceProvider"]
    checks = [
        {"key": "local-first", "label": "本地处理", "status": "ok", "detail": "源文件和输出都保存在本地项目目录，不上传云端。"},
        {"key": "editable-output", "label": "真实可编辑", "status": "ok" if job["outputMode"] == "pptx" else "warning", "detail": "PPTX 路线输出真实文字和形状；Web Deck 路线优先展示体验。"},
        {"key": "provider", "label": "模型 Provider", "status": "ok" if model_provider_ready(env, provider_config) else "warning", "detail": f"当前选择：{model_provider}。预览可运行，生产级生成会沿用该模型配置。"},
        {"key": "image-provider", "label": "图像 Provider", "status": "ok" if image_provider_ready(env, provider_config) else "warning", "detail": f"当前选择：{image_provider}。缺少 key 时会退回为本地预览和占位素材。"},
        {"key": "document-parser", "label": "资料解析", "status": "warning" if suffix in (DOCUMENT_SUFFIXES - {".docx"}) else "ok", "detail": "DOCX 已读取正文；PDF/XLSX/PPTX 已进入项目目录，完整解析由 Agent 工作流接管。" if suffix == ".docx" else "PDF/XLSX/PPTX 已进入项目目录；生产级全文解析由 Agent 工作流接管。" if suffix in DOCUMENT_SUFFIXES else "源资料已写入 source.md。"},
    ]
    if not env["python"]["bundledVenv"]:
        checks.append({"key": "venv", "label": "Python venv", "status": "missing", "detail": "未找到仓库 .venv；请先按 INSTALL.md 初始化依赖。"})
    if not env["optional"]["rust"]:
        checks.append({"key": "rust", "label": "Tauri 打包", "status": "warning", "detail": "未检测到 Rust/Cargo。Web 壳可用，原生 .app/.dmg 打包前需要安装 Rust。"})
    if provider_config.get("narrationEnabled"):
        checks.append({"key": "voice-provider", "label": "旁白语音", "status": "ok" if voice_provider_ready(env, provider_config) else "warning", "detail": f"当前选择：{voice_provider} / {provider_config.get('voiceId') or '未指定音色'}。"})
    return checks


def build_next_actions(project_path: Path, log_path: Path, generated_files: list[str], mode: str) -> list[dict[str, str]]:
    actions = [
        {
            "key": "open-result",
            "label": "打开结果文件",
            "detail": "先检查封面、结构和风格是否符合交付场景。",
            "path": generated_files[0] if generated_files else str(project_path),
        },
        {
            "key": "open-folder",
            "label": "打开项目文件夹",
            "detail": "查看 sources、outputs、preview、logs 和 manifest。",
            "path": str(project_path),
        },
        {
            "key": "agent-handoff",
            "label": "交给 Agent 深加工",
            "detail": "使用 sources/source.md 和 SKILL.md 进入生产级生成流程。",
            "path": str(project_path / "README.md"),
        },
        {
            "key": "formal-audit",
            "label": "正式商务审计",
            "detail": "运行 audit_formal_delivery.py、audit_design_completion.py 和 audit_visual_recipes.py，检查交付门槛、页面角色、页面配方和视觉完成度。",
            "path": str(project_path / "quality-report.json"),
        },
        {
            "key": "open-log",
            "label": "查看生成日志",
            "detail": "排查依赖、路径和导出问题。",
            "path": str(log_path),
        },
    ]
    if mode == "web":
        actions[0]["detail"] = "在浏览器中检查横向翻页、视觉节奏和分享效果。"
    return actions


def default_project_root(repo_root: Path) -> Path:
    return repo_root / DEFAULT_PROJECT_ROOT


def project_name_from_job(job: dict[str, Any]) -> str:
    source = job["source"]
    name = source["name"]
    if name:
        return sanitize_name(Path(name).stem)
    value = source["value"]
    if source["kind"] == "url":
        parsed = urlparse(value)
        return sanitize_name(parsed.netloc or "web-source")
    if source["kind"] == "file":
        return sanitize_name(Path(value).stem)
    first_line = value.splitlines()[0] if value.splitlines() else "desk"
    return sanitize_name(first_line[:32])


def create_project_dir(job: dict[str, Any], repo_root: Path) -> Path:
    base = Path(job["projectDir"]).expanduser() if job.get("projectDir") else default_project_root(repo_root)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    project_name = project_name_from_job(job)
    project_path = base / f"{project_name}-{job['outputMode']}-{timestamp}"
    for child in ("sources", "outputs", "logs", "preview"):
        (project_path / child).mkdir(parents=True, exist_ok=True)
    return project_path


def build_source_extraction(
    status: str,
    detail: str,
    generated_markdown_path: Path | None = None,
) -> dict[str, str]:
    payload = {
        "status": status,
        "detail": detail,
    }
    if generated_markdown_path is not None:
        payload["generatedMarkdownPath"] = str(generated_markdown_path)
    return payload


def convert_docx_to_markdown(repo_root: Path, input_path: Path, output_path: Path) -> str:
    converter_path = repo_root / "scripts" / "source_to_md" / "doc_to_md.py"
    if not converter_path.exists():
        raise RuntimeError("DOCX converter is missing at scripts/source_to_md/doc_to_md.py")

    spec = importlib.util.spec_from_file_location("ultimate_ppt_doc_to_md", converter_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("Unable to load DOCX converter.")

    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    convert_to_markdown = getattr(module, "convert_to_markdown", None)
    if not callable(convert_to_markdown):
        raise RuntimeError("DOCX converter does not expose convert_to_markdown().")

    stdout = io.StringIO()
    stderr = io.StringIO()
    with contextlib.redirect_stdout(stdout), contextlib.redirect_stderr(stderr):
        markdown = convert_to_markdown(str(input_path), str(output_path))

    if not markdown or not output_path.exists():
        converter_output = "\n".join(part for part in (stdout.getvalue().strip(), stderr.getvalue().strip()) if part)
        detail = converter_output or "DOCX converter returned no Markdown."
        raise RuntimeError(f"DOCX parsing failed: {detail}")
    return markdown


def convert_url_to_markdown(repo_root: Path, url: str, output_path: Path) -> str:
    converter_path = repo_root / "scripts" / "source_to_md" / "web_to_md.py"
    if not converter_path.exists():
        raise RuntimeError("URL converter is missing at scripts/source_to_md/web_to_md.py")

    spec = importlib.util.spec_from_file_location("ultimate_ppt_web_to_md", converter_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("Unable to load URL converter.")

    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    process_url = getattr(module, "process_url", None)
    if not callable(process_url):
        raise RuntimeError("URL converter does not expose process_url().")

    stdout = io.StringIO()
    stderr = io.StringIO()
    with contextlib.redirect_stdout(stdout), contextlib.redirect_stderr(stderr):
        ok, _source, error = process_url(url, str(output_path))

    if not ok or not output_path.exists():
        converter_output = "\n".join(part for part in (stdout.getvalue().strip(), stderr.getvalue().strip(), str(error or "").strip()) if part)
        detail = converter_output or "URL converter returned no Markdown."
        raise RuntimeError(f"URL parsing failed: {detail}")

    markdown = output_path.read_text(encoding="utf-8", errors="replace")
    if not markdown.strip():
        raise RuntimeError("URL parsing failed: URL converter wrote an empty Markdown file.")
    return markdown


def source_to_text(job: dict[str, Any], project_path: Path, repo_root: Path) -> tuple[str, str, dict[str, str]]:
    source = job["source"]
    kind = source["kind"]
    value = source["value"]
    sources_dir = project_path / "sources"
    generated_markdown_path = sources_dir / "source.md"

    if kind in {"text", "markdown"}:
        filename = "source.md" if kind == "markdown" else "source.txt"
        (sources_dir / filename).write_text(value, encoding="utf-8")
        if filename != "source.md":
            generated_markdown_path.write_text(value, encoding="utf-8")
        return (
            value,
            filename,
            build_source_extraction(
                "extracted",
                "已读取粘贴内容并写入本地项目。",
                generated_markdown_path,
            ),
        )

    if kind == "url":
        (sources_dir / "source.url").write_text(value + "\n", encoding="utf-8")
        try:
            text = convert_url_to_markdown(repo_root, value, generated_markdown_path)
            return (
                text,
                "source.url",
                build_source_extraction(
                    "extracted",
                    "已抓取 URL 正文并生成可用于 PPTX/Web 的 source.md。",
                    generated_markdown_path,
                ),
            )
        except Exception as err:
            text = (
                f"URL source: {value}\n\n"
                "Desktop URL extraction failed, so this project was staged for Agent handoff.\n\n"
                f"Error: {err}"
            )
            generated_markdown_path.write_text(text, encoding="utf-8")
            return (
                text,
                "source.url",
                build_source_extraction(
                    "handoffRequired",
                    f"已保存 URL；桌面抓取失败，网页抓取和事实校验由 Agent 工作流接管。错误：{clip_text(err, 120)}",
                    generated_markdown_path,
                ),
            )

    source_path = Path(value).expanduser()
    if not source_path.exists():
        raise FileNotFoundError(f"Source file not found: {source_path}")
    copied_path = sources_dir / source_path.name
    if source_path.resolve() != copied_path.resolve():
        shutil.copy2(source_path, copied_path)

    if source_path.suffix.lower() in TEXT_SUFFIXES:
        text = source_path.read_text(encoding="utf-8", errors="replace")
        generated_markdown_path.write_text(text, encoding="utf-8")
        extraction = build_source_extraction(
            "extracted",
            "已读取文本类文件内容并写入 source.md。",
            generated_markdown_path,
        )
    elif source_path.suffix.lower() == ".docx":
        text = convert_docx_to_markdown(repo_root, copied_path, generated_markdown_path)
        extraction = build_source_extraction(
            "extracted",
            "已解析 DOCX 正文并生成可用于 PPTX/Web 的 source.md。",
            generated_markdown_path,
        )
    else:
        text = (
            f"Imported file: {source_path.name}\n\n"
            "This desktop draft created a local project shell. For full document extraction, "
            "run the full ultimate-ppt-master workflow from SKILL.md."
        )
        generated_markdown_path.write_text(text, encoding="utf-8")
        status = "handoffRequired" if source_path.suffix.lower() in DOCUMENT_SUFFIXES else "copied"
        detail = "已复制源文件；该格式的完整解析由 Agent 工作流接管。" if status == "handoffRequired" else "已复制源文件并创建项目壳。"
        extraction = build_source_extraction(status, detail, generated_markdown_path)
    return text, source_path.name, extraction


def looks_chinese(text: str) -> bool:
    return bool(re.search(r"[\u4e00-\u9fff]", text))


def extract_lines(text: str) -> list[str]:
    cleaned: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        line = re.sub(r"^#{1,6}\s*", "", line)
        line = re.sub(r"^[-*+]\s*", "", line)
        line = re.sub(r"^\d+[.)]\s*", "", line)
        if line and line not in cleaned:
            cleaned.append(line)
    if not cleaned:
        cleaned = ["Untitled presentation", "Add source material to generate a richer deck."]
    return cleaned[:60]


def is_section_heading(line: str) -> bool:
    return bool(re.match(r"^([一二三四五六七八九十]+、|\d+[、.．])", line))


def clip_text(value: Any, limit: int = 88) -> str:
    text = str(value).strip()
    return text if len(text) <= limit else f"{text[:limit - 1]}…"


def clip_sentence(value: Any, limit: int = 58) -> str:
    text = re.sub(r"\s+", " ", str(value).strip())
    text = re.sub(r"^[_*`#>-]+|[_*`#>-]+$", "", text).strip()
    return text if len(text) <= limit else f"{text[:limit - 1]}…"


def split_title(value: str, line_limit: int = 16, max_lines: int = 3) -> list[str]:
    title = re.sub(r"\s+", "", value.strip())
    if not title:
        return ["Untitled"]
    lines = [title[idx: idx + line_limit] for idx in range(0, len(title), line_limit)]
    closing = "，。！？；：、）》」』”’"
    opening = "《（「『“‘"
    for idx in range(1, len(lines)):
        if lines[idx] and lines[idx][0] in closing and lines[idx - 1]:
            lines[idx - 1] += lines[idx][0]
            lines[idx] = lines[idx][1:]
        if lines[idx - 1] and lines[idx - 1][-1] in opening and lines[idx]:
            lines[idx] = lines[idx - 1][-1] + lines[idx]
            lines[idx - 1] = lines[idx - 1][:-1]
    lines = [line for line in lines if line]
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        lines[-1] = lines[-1].rstrip("…") + "…"
    return lines


def cover_title_lines(value: str, line_limit: int = 12, max_lines: int = 3) -> list[str]:
    raw = re.sub(r"\s+", "", value.strip())
    quoted = re.search(r"[“\"]([^”\"]{2,16})[”\"]", raw)
    if quoted:
        lead = quoted.group(1)
        tail = raw[quoted.end():]
        tail = re.sub(r"^[暨及和与·:：-]+", "", tail)
        tail = re.sub(r"(的)?(请示|方案|材料|初稿|汇报|报告|通知).*$", "", tail)
        tail = tail.replace("2026年", "2026")
        if max_lines >= 3 and tail.endswith("消费主题活动") and len(tail) > len("消费主题活动"):
            tail_lines = [tail[: -len("消费主题活动")], "消费主题活动"]
        else:
            tail_lines = split_title(tail, line_limit, max(1, max_lines - 1)) if tail else []
        return [lead, *tail_lines][:max_lines]
    raw = re.sub(r"^关于(举办|开展|推进|落实)", "", raw)
    raw = re.sub(r"(的)?(请示|方案|材料|初稿|汇报|报告|通知).*$", "", raw)
    return split_title(raw, line_limit, max_lines)


def slide_title(value: str, limit: int = 34) -> str:
    return clip_sentence(value, limit)


def slide_body(value: str, limit: int = 50) -> str:
    return clip_sentence(value, limit)


def build_sections(lines: list[str]) -> list[dict[str, Any]]:
    sections: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for line in lines[1:]:
        if is_section_heading(line):
            current = {"title": clip_text(line, 56), "bullets": []}
            sections.append(current)
            continue
        if current is None:
            current = {"title": "核心信息", "bullets": []}
            sections.append(current)
        if len(current["bullets"]) < 5:
            current["bullets"].append(clip_text(line))
    return [section for section in sections if section["bullets"]]


def build_outline(text: str, style: str, mode: str) -> list[dict[str, Any]]:
    lines = extract_lines(text)
    chinese = looks_chinese(text)
    title = lines[0][:80]
    if chinese:
        headings = {
            "context": "核心信息",
            "workflow": "生成路径",
            "output": "交付结果",
            "next": "下一步",
        }
        mode_label = "可编辑 PPTX" if mode == "pptx" else "杂志风网页 PPT"
        workflow_bullets = ["导入资料", "解析 source.md", "锁定结构与风格", "生成预览", "校验并导出"]
        output_bullets = [mode_label, "本地项目文件夹", "预览产物", "Agent 生产级工作流说明"]
        next_bullets = ["打开生成文件", "检查页面节奏", "使用完整 skill 继续生产级打磨"]
    else:
        headings = {
            "context": "Core Signal",
            "workflow": "Generation Path",
            "output": "Delivery Output",
            "next": "Next Step",
        }
        mode_label = "Editable PPTX" if mode == "pptx" else "Magazine Web Deck"
        workflow_bullets = ["Import source", "Extract source.md", "Lock structure and style", "Generate preview", "Verify and export"]
        output_bullets = [mode_label, "Local project folder", "Preview artifact", "Runbook for full agent workflow"]
        next_bullets = ["Open the generated files", "Review page rhythm", "Continue with the full skill for production quality"]

    body = lines[1:] or lines
    outline = [
        {"title": title, "eyebrow": "Ultimate PPT Master", "bullets": [mode_label, f"Style preset: {style}"]},
        {"title": headings["context"], "eyebrow": "01", "bullets": [clip_text(item) for item in body[:4]]},
    ]
    target = 10 if mode == "pptx" and len(lines) >= 14 else 8 if mode == "web" and len(lines) >= 14 else 6
    content_limit = max(2, target - 3)
    used_titles = {slide["title"] for slide in outline}

    for section in build_sections(lines):
        if len(outline) >= content_limit:
            break
        if section["title"] in used_titles:
            continue
        outline.append(
            {
                "title": section["title"],
                "eyebrow": f"{len(outline):02d}",
                "bullets": section["bullets"][:5],
            }
        )
        used_titles.add(section["title"])

    body_index = 0
    fallback_titles = ["关键安排", "资源协同", "执行重点", "预算与保障", "风险与下一步"] if chinese else ["Key Moves", "Resource Fit", "Execution Focus", "Budget and Support", "Risks and Next"]
    while len(outline) < content_limit and body_index < len(body):
        chunk = [clip_text(item) for item in body[body_index:body_index + 4]]
        body_index += 4
        if not chunk:
            break
        fallback_title = fallback_titles[(len(outline) - 2) % len(fallback_titles)]
        outline.append({"title": fallback_title, "eyebrow": f"{len(outline):02d}", "bullets": chunk})

    outline.append({"title": headings["workflow"], "eyebrow": f"{len(outline):02d}", "bullets": workflow_bullets})
    outline.append({"title": headings["output"], "eyebrow": f"{len(outline):02d}", "bullets": output_bullets})
    outline.append({"title": headings["next"], "eyebrow": f"{len(outline):02d}", "bullets": next_bullets})
    return outline


def write_preview_svg(outline: list[dict[str, Any]], project_path: Path, style: str) -> str:
    colors = {
        "business": ("#FFF8EC", "#172033", "#F97316"),
        "consulting": ("#F7F5EF", "#111827", "#2563EB"),
        "academic": ("#F4F8FF", "#172033", "#10B981"),
        "editorial": ("#FFF8EC", "#2F241B", "#C2410C"),
        "swiss": ("#F7FAFC", "#111827", "#E11D48"),
    }
    bg, ink, accent = colors.get(style, colors["business"])
    slide = outline[0]
    bullet_text = "".join(
        f'<text x="88" y="{430 + idx * 48}" fill="{ink}" font-family="Avenir Next, PingFang SC, sans-serif" font-size="30">{html.escape(str(item))}</text>'
        for idx, item in enumerate(slide["bullets"][:4])
    )
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1600 900">
  <rect width="1600" height="900" fill="{bg}"/>
  <path d="M-40 190C260 40 480 180 720 120C980 56 1120 194 1640 104" fill="none" stroke="{accent}" stroke-width="44" opacity="0.16"/>
  <rect x="70" y="70" width="1460" height="760" rx="44" fill="#FFFFFF" opacity="0.82"/>
  <text x="88" y="154" fill="{accent}" font-family="Avenir Next, PingFang SC, sans-serif" font-size="28" font-weight="800">{html.escape(slide['eyebrow'])}</text>
  <text x="88" y="292" fill="{ink}" font-family="Avenir Next, PingFang SC, sans-serif" font-size="72" font-weight="850">{html.escape(slide['title'])}</text>
  {bullet_text}
  <text x="88" y="750" fill="#667085" font-family="Avenir Next, PingFang SC, sans-serif" font-size="24">Production draft · Full refinement workflow remains agent-driven</text>
</svg>"""
    path = project_path / "preview" / "cover.svg"
    path.write_text(svg, encoding="utf-8")
    return svg


def shape_rect() -> Any:
    return MSO_SHAPE.RECTANGLE if MSO_SHAPE is not None else 1


def add_text_box(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    align: Any | None = None,
    font_name: str = "Aptos",
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    p = frame.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    return box


def add_filled_rect(slide: Any, left: float, top: float, width: float, height: float, fill: tuple[int, int, int], line: tuple[int, int, int] | None = None) -> Any:
    shape = slide.shapes.add_shape(shape_rect(), Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = RGBColor(*line)
    return shape


def ppt_palette(style: str) -> dict[str, tuple[int, int, int]]:
    palettes = {
        "business": {
            "paper": (248, 246, 240),
            "ink": (18, 28, 45),
            "muted": (85, 95, 112),
            "accent": (231, 84, 41),
            "accent2": (31, 97, 141),
            "soft": (255, 238, 221),
            "card": (255, 255, 255),
        },
        "consulting": {
            "paper": (246, 248, 251),
            "ink": (17, 24, 39),
            "muted": (75, 85, 99),
            "accent": (37, 99, 235),
            "accent2": (15, 118, 110),
            "soft": (226, 232, 240),
            "card": (255, 255, 255),
        },
        "academic": {
            "paper": (244, 248, 255),
            "ink": (23, 43, 77),
            "muted": (91, 105, 135),
            "accent": (16, 122, 105),
            "accent2": (79, 70, 229),
            "soft": (222, 241, 236),
            "card": (255, 255, 255),
        },
        "editorial": {
            "paper": (241, 239, 234),
            "ink": (23, 21, 19),
            "muted": (98, 89, 79),
            "accent": (194, 65, 12),
            "accent2": (62, 88, 121),
            "soft": (232, 229, 222),
            "card": (255, 252, 246),
        },
        "swiss": {
            "paper": (255, 255, 255),
            "ink": (10, 10, 10),
            "muted": (82, 82, 82),
            "accent": (0, 86, 255),
            "accent2": (225, 29, 72),
            "soft": (236, 242, 255),
            "card": (246, 247, 249),
        },
    }
    return palettes.get(style, palettes["business"])


def generate_pptx(outline: list[dict[str, Any]], project_path: Path, style: str) -> Path:
    if Presentation is None:
        raise RuntimeError("python-pptx is not available. Install requirements.txt first.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    palette = ppt_palette(style)
    paper = palette["paper"]
    ink = palette["ink"]
    muted = palette["muted"]
    accent = palette["accent"]
    accent2 = palette["accent2"]
    soft = palette["soft"]
    card = palette["card"]
    contracts = page_contract_for_outline(outline)

    for idx, item in enumerate(outline):
        contract = contracts[idx]
        role = contract["page_role"]
        recipe = contract["page_recipe_id"]
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*paper)

        add_filled_rect(slide, 0, 0, 0.16, 7.5, accent)
        add_text_box(slide, 0.48, 0.34, 3.1, 0.28, f"{item['eyebrow']} / {idx + 1:02d}", 9, accent, True, font_name="Aptos Mono")
        add_text_box(slide, 11.2, 0.34, 1.45, 0.28, "DRAFT", 9, muted, True, PP_ALIGN.RIGHT if PP_ALIGN else None, "Aptos Mono")

        if idx == 0:
            add_filled_rect(slide, 8.65, 0, 4.68, 7.5, accent)
            add_filled_rect(slide, 0.62, 5.78, 2.6, 0.08, accent2)
            title_lines = cover_title_lines(str(item["title"]), 14, 3)
            y = 1.18
            for line in title_lines:
                add_text_box(slide, 0.62, y, 7.65, 0.7, line, 34, ink, True, font_name="Aptos Display")
                y += 0.72
            lead = " / ".join(slide_body(bullet, 26) for bullet in item["bullets"][:2])
            add_text_box(slide, 0.66, 4.2, 6.8, 0.76, lead, 18, muted)
            add_text_box(slide, 9.04, 1.05, 3.38, 0.55, "Editable PowerPoint", 24, (255, 255, 255), True)
            add_text_box(slide, 9.08, 1.82, 3.25, 1.5, "真实文字、形状、结构和本地项目链路。桌面端先给出可审阅草稿，完整精修交给 Agent 工作流。", 15, (255, 245, 238))
            add_text_box(slide, 9.08, 6.28, 3.25, 0.32, "Ultimate PPT Master · hybrid-editable", 10, (255, 245, 238), font_name="Aptos Mono")
            continue

        if role == "closing":
            bg.fore_color.rgb = RGBColor(*ink)
            add_filled_rect(slide, 0, 0, 13.333, 0.18, accent)
            add_text_box(slide, 0.72, 0.62, 1.7, 0.4, f"{idx + 1:02d}", 26, accent, True, font_name="Aptos Display")
            add_text_box(slide, 0.72, 2.04, 10.8, 1.25, slide_title(str(item["title"]), 28), 38, (255, 255, 255), True, font_name="Aptos Display")
            add_text_box(slide, 0.76, 4.0, 8.8, 0.9, " / ".join(slide_body(b, 22) for b in item["bullets"][:3]), 17, (221, 226, 235))
            add_text_box(slide, 9.5, 5.85, 2.55, 0.3, "HYBRID EDITABLE", 10, (221, 226, 235), True, PP_ALIGN.RIGHT if PP_ALIGN else None, "Aptos Mono")
            continue

        add_text_box(slide, 0.68, 0.96, 11.2, 0.72, slide_title(str(item["title"])), 27, ink, True, font_name="Aptos Display")
        bullets = [slide_body(bullet) for bullet in item["bullets"][:6]]

        if recipe == "process_flow.horizontal_steps":
            add_filled_rect(slide, 0.72, 3.08, 11.3, 0.04, accent)
            for bullet_idx, bullet in enumerate(bullets[:5]):
                left = 0.86 + bullet_idx * 2.22
                add_filled_rect(slide, left, 2.38, 1.72, 1.5, card, soft)
                add_text_box(slide, left + 0.16, 2.58, 0.5, 0.28, f"{bullet_idx + 1:02d}", 13, accent, True, font_name="Aptos Mono")
                add_text_box(slide, left + 0.16, 3.0, 1.32, 0.54, bullet, 12, ink, bullet_idx == 0)
            add_text_box(slide, 0.78, 5.42, 9.2, 0.5, "流程页采用可编辑节点，ChatGPT 只提供无文字流程氛围层。", 13, muted)
        elif recipe == "metric_panel.large_number_strip":
            for bullet_idx, bullet in enumerate(bullets[:4]):
                left = 0.74 + bullet_idx * 2.92
                add_filled_rect(slide, left, 2.18, 2.42, 2.28, card, soft)
                add_text_box(slide, left + 0.22, 2.5, 0.72, 0.44, f"{bullet_idx + 1}", 28, accent, True, font_name="Aptos Display")
                add_text_box(slide, left + 0.22, 3.34, 1.84, 0.72, bullet, 13, ink)
            add_filled_rect(slide, 0.74, 5.32, 8.8, 0.08, accent2)
            add_text_box(slide, 0.74, 5.58, 8.2, 0.48, "数字、金额和单位必须保留为可编辑文本。", 13, muted)
        elif recipe == "comparison_matrix.two_column_delta":
            add_filled_rect(slide, 0.78, 2.02, 5.25, 3.9, card, soft)
            add_filled_rect(slide, 6.42, 2.02, 5.25, 3.9, card, soft)
            add_text_box(slide, 1.08, 2.34, 2.8, 0.34, "原状态 / Before", 15, accent, True)
            add_text_box(slide, 6.72, 2.34, 2.8, 0.34, "升级后 / After", 15, accent2, True)
            for bullet_idx, bullet in enumerate(bullets[:3]):
                add_text_box(slide, 1.08, 3.02 + bullet_idx * 0.72, 4.1, 0.42, bullet, 14, ink)
            for bullet_idx, bullet in enumerate((bullets[3:6] or bullets[:3])):
                add_text_box(slide, 6.72, 3.02 + bullet_idx * 0.72, 4.1, 0.42, bullet, 14, ink)
        elif recipe == "risk_callout.qa_stack":
            add_filled_rect(slide, 0.78, 2.02, 10.9, 0.56, soft)
            add_text_box(slide, 1.04, 2.15, 6.2, 0.28, "办理条件 / 风险边界 / 客户疑问", 14, accent, True)
            for bullet_idx, bullet in enumerate(bullets[:4]):
                top = 2.92 + bullet_idx * 0.78
                add_text_box(slide, 0.98, top, 0.52, 0.3, "Q", 15, accent, True, font_name="Aptos Mono")
                add_text_box(slide, 1.48, top, 8.6, 0.35, bullet, 14, ink)
        elif recipe == "action_roadmap.owner_timeline":
            for bullet_idx, bullet in enumerate(bullets[:5]):
                top = 2.08 + bullet_idx * 0.76
                add_filled_rect(slide, 0.84, top + 0.12, 0.28, 0.28, accent)
                add_text_box(slide, 1.38, top, 1.2, 0.3, f"T+{bullet_idx}", 12, accent, True, font_name="Aptos Mono")
                add_text_box(slide, 2.58, top, 7.8, 0.36, bullet, 14, ink)
            add_filled_rect(slide, 0.96, 2.32, 0.04, 3.28, soft)
        elif recipe == "evidence_board.source_table":
            add_filled_rect(slide, 0.76, 2.02, 10.8, 0.44, soft)
            add_text_box(slide, 1.0, 2.13, 2.0, 0.22, "证据", 11, accent, True)
            add_text_box(slide, 5.6, 2.13, 2.0, 0.22, "口径", 11, accent, True)
            for bullet_idx, bullet in enumerate(bullets[:5]):
                top = 2.62 + bullet_idx * 0.68
                add_filled_rect(slide, 0.76, top, 10.8, 0.48, card, soft)
                add_text_box(slide, 1.0, top + 0.12, 0.5, 0.22, f"{bullet_idx + 1:02d}", 10, accent, True, font_name="Aptos Mono")
                add_text_box(slide, 1.58, top + 0.08, 8.6, 0.28, bullet, 12, ink)
        else:
            add_filled_rect(slide, 0.72, 2.3, 7.2, 2.7, card)
            add_filled_rect(slide, 0.72, 2.3, 0.12, 2.7, accent)
            for bullet_idx, bullet in enumerate(bullets[:4]):
                add_text_box(slide, 1.06, 2.64 + bullet_idx * 0.58, 6.28, 0.36, bullet, 15, ink, bullet_idx == 0)
            add_filled_rect(slide, 8.6, 2.28, 3.42, 2.72, soft)
            add_text_box(slide, 9.0, 2.64, 2.62, 0.45, "页面配方", 18, accent, True)
            add_text_box(slide, 9.0, 3.24, 2.52, 1.0, recipe, 12, muted, font_name="Aptos Mono")

        add_filled_rect(slide, 0.72, 6.72, 9.5, 0.02, soft)
        add_text_box(slide, 0.72, 6.86, 7.4, 0.26, f"Hybrid editable · {recipe}", 8, muted, font_name="Aptos Mono")
        add_text_box(slide, 10.8, 6.82, 1.4, 0.28, f"{idx + 1:02d}", 10, accent, True, PP_ALIGN.RIGHT if PP_ALIGN else None, "Aptos Mono")
    output = project_path / "outputs" / "ultimate-ppt-master-preview.pptx"
    prs.save(output)
    return output


def copy_web_assets(project_path: Path, repo_root: Path) -> None:
    assets_dir = project_path / "outputs" / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    motion_src = repo_root / "assets" / "magazine-web" / "motion.min.js"
    if motion_src.exists():
        shutil.copy2(motion_src, assets_dir / "motion.min.js")


def replace_template_slides(template: str, slides_html: str) -> str:
    if "<!-- SLIDES_HERE -->" in template:
        return template.replace("<!-- SLIDES_HERE -->", slides_html)
    marker = "<!-- SLIDES_HERE"
    if marker not in template:
        raise RuntimeError("Magazine Web template is missing SLIDES_HERE marker.")
    start = template.index(marker)
    end_marker = "\n</div>\n\n<div id=\"nav\""
    end = template.find(end_marker, start)
    if end < 0:
        raise RuntimeError("Unable to locate deck closing tag in Magazine Web template.")
    return template[:start] + slides_html + template[end:]


def chrome(title: str, page: int, total: int) -> str:
    return f"""<div class="chrome">
    <div>{html.escape(title)} · Ultimate PPT Master</div>
    <div>{page:02d} / {total:02d}</div>
  </div>"""


def foot(label: str) -> str:
    return f"""<div class="foot">
    <div>{html.escape(label)}</div>
    <div>— · —</div>
  </div>"""


def build_editorial_sections(outline: list[dict[str, Any]], deck_title: str) -> str:
    total = len(outline)
    sections: list[str] = []
    contracts = page_contract_for_outline(outline)
    title_lines = cover_title_lines(str(outline[0]["title"]), 10, 3)
    title = "<br>".join(html.escape(line) for line in title_lines)
    title_size = "min(4.5vw,8.2vh)" if len(str(outline[0]["title"])) > 24 else "min(6.4vw,11vh)"
    lead_items = outline[1]["bullets"][:2] if len(outline) > 1 else outline[0]["bullets"][:2]
    subtitle = " · ".join(html.escape(slide_body(item, 30)) for item in lead_items)
    cover_points = []
    for idx, item in enumerate(outline[1:4], start=1):
        point = slide_body(item["bullets"][0] if item["bullets"] else item["title"], 34)
        cover_points.append(
            f"""<div class="stat-card" data-anim>
        <div class="stat-label">Theme {idx:02d}</div>
        <div class="stat-nb">{idx}<span class="stat-unit">项</span></div>
        <div class="stat-note">{html.escape(point)}</div>
      </div>"""
        )
    cover_cards = "\n".join(cover_points)
    sections.append(f"""<section class="slide hero light" data-layout="{html.escape(contracts[0]['page_recipe_id'])}" data-animate="hero">
  {chrome(deck_title, 1, total)}
  <div class="frame" style="display:grid;grid-template-columns:1.28fr .82fr;gap:5vw;align-items:center;min-height:80vh">
    <div style="display:grid;gap:3vh;align-content:center">
      <div class="kicker" data-anim>Production Draft · Local First</div>
      <h1 class="h-hero" style="font-size:{title_size};line-height:1.08;max-width:58vw" data-anim>{title}</h1>
      <p class="lead" style="max-width:50vw" data-anim>{subtitle}</p>
      <div class="meta-row" data-anim><span>Editable PPTX</span><span>·</span><span>Magazine Web Deck</span><span>·</span><span>Agent Ready</span></div>
    </div>
    <div style="display:grid;gap:2.4vh;align-content:center">
      <div class="callout" style="border-left-color:#c2410c;background:rgba(194,65,12,.08)" data-anim>
        <div class="q-big">从原始办公材料进入桌面生成链路，先给出可审阅、可复用、可继续精修的演示草稿。</div>
        <span class="cite">DOCX → Markdown → Deck</span>
      </div>
      <div class="grid-3" style="gap:2vw;flex:initial">{cover_cards}</div>
    </div>
  </div>
  {foot("真实资料 · 本地生成 · 可继续深加工")}
</section>""")

    for idx, slide in enumerate(outline[1:], start=2):
        contract = contracts[idx - 1]
        recipe = contract["page_recipe_id"]
        bullets = [html.escape(slide_body(item, 48)) for item in slide["bullets"][:6]]
        title_html = html.escape(slide_title(str(slide["title"]), 32))
        if recipe == "metric_panel.large_number_strip":
            cards = "\n".join(
                f"""<div class="stat-card" data-anim>
        <div class="stat-label">Signal {n + 1:02d}</div>
        <div class="stat-nb">{n + 1}<span class="stat-unit">项</span></div>
        <div class="stat-note">{bullet}</div>
      </div>"""
                for n, bullet in enumerate(bullets[:6])
            )
            sections.append(f"""<section class="slide light" data-layout="{html.escape(recipe)}">
  {chrome(deck_title, idx, total)}
  <div class="frame" style="padding-top:5vh">
    <div class="kicker" data-anim>{html.escape(str(slide['eyebrow']))}</div>
    <h2 class="h-xl" data-anim>{title_html}</h2>
    <div class="grid-6" style="margin-top:5vh">{cards}</div>
  </div>
  {foot("核心信号")}
</section>""")
        elif recipe in {"comparison_matrix.two_column_delta", "risk_callout.qa_stack"}:
            left = bullets[:3]
            right = bullets[3:6] or bullets[:3]
            left_html = "<br>".join(left)
            right_html = "".join(f"<div class=\"rowline\" data-anim><div class=\"k\">{n + 1:02d}</div><div class=\"v\">{bullet}</div><div class=\"m\">Action</div></div>" for n, bullet in enumerate(right))
            sections.append(f"""<section class="slide dark" data-layout="{html.escape(recipe)}" data-animate="directional">
  {chrome(deck_title, idx, total)}
  <div class="frame grid-2-7-5" style="padding-top:6vh">
    <div>
      <div class="kicker" data-anim>{html.escape(str(slide['eyebrow']))}</div>
      <h2 class="h-xl" style="font-size:6.8vw" data-anim>{title_html}</h2>
      <div class="callout" style="margin-top:4vh" data-anim>{left_html}</div>
    </div>
    <div style="padding-top:7vh">{right_html}</div>
  </div>
  {foot("结构化安排")}
</section>""")
        elif recipe in {"process_flow.horizontal_steps", "action_roadmap.owner_timeline"}:
            steps = "\n".join(
                f"""<div class="step" data-anim="step">
          <div class="step-nb">{n + 1:02d}</div>
          <div class="step-title">{html.escape(slide_title(bullet, 12))}</div>
          <div class="step-desc">{bullet}</div>
        </div>"""
                for n, bullet in enumerate(bullets[:5])
            )
            sections.append(f"""<section class="slide light" data-layout="{html.escape(recipe)}" data-animate="pipeline">
  {chrome(deck_title, idx, total)}
  <div class="frame">
    <div class="kicker">Pipeline · 推进路径</div>
    <h2 class="h-xl">{title_html}</h2>
    <div class="pipeline-section">
      <div class="pipeline-label">执行链路 · Action Flow</div>
      <div class="pipeline">{steps}</div>
    </div>
  </div>
  {foot("流程与责任")}
</section>""")
        else:
            body = " ".join(bullets[:3])
            sections.append(f"""<section class="slide hero {'light' if idx % 2 else 'dark'}" data-layout="{html.escape(recipe)}" data-animate="quote">
  {chrome(deck_title, idx, total)}
  <div class="frame" style="display:grid;gap:5vh;align-content:center;min-height:80vh">
    <div class="kicker" data-anim>{html.escape(str(slide['eyebrow']))}</div>
    <h2 class="h-xl" data-anim>{title_html}</h2>
    <p class="lead" style="max-width:58vw" data-anim>{html.escape(body)}</p>
  </div>
  {foot("观点页")}
</section>""")
    return "\n\n".join(sections)


def build_swiss_sections(outline: list[dict[str, Any]], deck_title: str) -> str:
    total = len(outline)
    sections: list[str] = []
    contracts = page_contract_for_outline(outline)
    title = "<br>".join(html.escape(line) for line in cover_title_lines(str(outline[0]["title"]), 10, 3))
    lead_items = outline[1]["bullets"][:2] if len(outline) > 1 else outline[0]["bullets"][:2]
    lead = " · ".join(html.escape(slide_body(item, 30)) for item in lead_items)
    sections.append(f"""<section class="slide accent" data-layout="{html.escape(contracts[0]['page_recipe_id'])}" data-animate="hero">
  <div class="canvas-card">
    <canvas class="ascii-bg" aria-hidden="true"></canvas>
    <div class="chrome-min"><div class="l">{html.escape(deck_title)} · Field Note</div><div class="r">SS · 01 / {total:02d}</div></div>
    <div style="flex:1;padding:0;display:grid;grid-template-rows:auto 1fr auto;gap:2.6vh">
      <div data-anim="kicker" class="t-meta" style="color:rgba(255,255,255,.78);letter-spacing:.22em">LOCAL-FIRST PRESENTATION</div>
      <h1 data-anim="title" style="align-self:center;font-family:var(--sans),var(--sans-zh);font-weight:200;font-size:min(6.2vw,11vh);line-height:1.02;letter-spacing:-.025em;color:#fff">{title}</h1>
      <div data-anim="bottom" style="display:grid;grid-template-rows:auto auto;gap:1.6vh;border-top:1px solid rgba(255,255,255,.22);padding-top:2vh">
        <div class="lead" style="max-width:54ch;color:rgba(255,255,255,.86);font-weight:300">{lead}</div>
        <div style="display:flex;justify-content:space-between;align-items:end"><div class="t-meta" style="color:rgba(255,255,255,.6)">Ultimate PPT Master</div><div class="t-meta" style="color:rgba(255,255,255,.6)">→ swipe / arrow keys</div></div>
      </div>
    </div>
  </div>
</section>""")

    layouts = ["S02", "S08", "S14", "S19", "S11", "S16", "S03"]
    for idx, slide in enumerate(outline[1:], start=2):
        contract = contracts[idx - 1]
        recipe = contract["page_recipe_id"]
        layout = layouts[(idx - 2) % len(layouts)]
        bullets = [html.escape(slide_body(item, 48)) for item in slide["bullets"][:6]]
        title_html = html.escape(slide_title(str(slide["title"]), 28))
        if layout == "S02":
            cards = "\n".join(
                f"""<div class="stat-card {'accent-top' if n == 0 else 'thin'}" data-anim>
          <div class="stat-label">Point {n + 1:02d}</div>
          <div class="stat-nb">{n + 1}<span class="stat-unit">项</span></div>
          <div class="stat-note">{bullet}</div>
        </div>"""
                for n, bullet in enumerate(bullets[:6])
            )
            body = f"""<div class="grid-6" style="margin-top:4vh">{cards}</div>"""
        elif layout == "S14":
            blocks = "\n".join(
                f"""<div class="stack-block {'b-accent' if n == 1 else 'b-grey'}" data-anim>
          <div class="layer-nb">L{n + 1:02d}</div>
          <div class="layer-ttl">{html.escape(slide_title(bullet, 14))}</div>
          <div class="layer-desc">{bullet}</div>
        </div>"""
                for n, bullet in enumerate(bullets[:4])
            )
            body = f"""<div style="display:grid;grid-template-columns:repeat({max(1, min(4, len(bullets)))},1fr);gap:16px;margin-top:4vh">{blocks}</div>"""
        else:
            cards = "\n".join(
                f"""<div class="sub-card {'accent' if n == 0 and idx % 2 == 0 else ''}" data-anim>
          <div class="nb-corner">{n + 1:02d}</div>
          <div class="ttl">{html.escape(slide_title(bullet, 14))}</div>
          <div class="desc">{bullet}</div>
        </div>"""
                for n, bullet in enumerate(bullets[:4])
            )
            body = f"""<div style="display:grid;grid-template-columns:repeat(2,1fr);gap:16px;margin-top:4vh">{cards}</div>"""
        sections.append(f"""<section class="slide {'dark' if idx % 3 == 0 else 'grey'}" data-layout="{html.escape(recipe)}" data-template-layout="{layout}" data-animate="grid-reveal">
  <div class="canvas-card">
    <div class="chrome-min"><div class="l">{html.escape(deck_title)}</div><div class="r">{idx:02d} / {total:02d}</div></div>
    <div style="flex:1;padding:0;display:grid;grid-template-rows:auto 1fr auto;gap:3vh">
      <div data-anim="line" style="display:flex;flex-direction:column;gap:1.2vh">
        <div class="t-meta">{html.escape(str(slide['eyebrow']))}</div>
        <h2 class="h-xl-zh" style="font-size:min(5.6vw,10vh);line-height:1.02">{title_html}</h2>
      </div>
      <div data-anim="up">{body}</div>
      <div class="t-meta" style="text-align:right;color:var(--text-helper)">Production draft · Agent ready</div>
    </div>
  </div>
</section>""")
    return "\n\n".join(sections)


def generate_web_deck(outline: list[dict[str, Any]], project_path: Path, style: str, repo_root: Path) -> tuple[Path, str]:
    copy_web_assets(project_path, repo_root)
    template_name = "template-swiss.html" if style == "swiss" else "template.html"
    template_path = repo_root / "assets" / "magazine-web" / template_name
    template = template_path.read_text(encoding="utf-8")
    deck_title = slide_title(str(outline[0]["title"]), 22)
    sections = build_swiss_sections(outline, deck_title) if style == "swiss" else build_editorial_sections(outline, deck_title)
    deck = replace_template_slides(template, sections)
    deck = deck.replace("[必填] 替换为 PPT 标题 · Deck Title", f"{deck_title} · Ultimate PPT Master")
    deck = re.sub(r"\[必填\][^<\n]*", "", deck)
    output = project_path / "outputs" / "index.html"
    output.write_text(deck, encoding="utf-8")
    return output, deck


def write_narration_handoff(project_path: Path, outline: list[dict[str, Any]], provider_config: dict[str, Any]) -> list[Path]:
    if not provider_config.get("narrationEnabled") or provider_config.get("voiceProvider") == "none":
        return []

    notes_dir = project_path / "notes"
    audio_dir = project_path / "audio"
    notes_dir.mkdir(exist_ok=True)
    audio_dir.mkdir(exist_ok=True)

    for idx, slide in enumerate(outline, start=1):
        bullets = "\n".join(f"- {item}" for item in slide["bullets"][:5])
        note = f"# Slide {idx}: {slide['title']}\n\n{bullets}\n"
        (notes_dir / f"slide-{idx:02d}.md").write_text(note, encoding="utf-8")

    settings = {
        "enabled": True,
        "provider": provider_config.get("voiceProvider"),
        "voiceId": provider_config.get("voiceId"),
        "rate": provider_config.get("voiceRate") or "+0%",
        "notesDir": str(notes_dir),
        "audioDir": str(audio_dir),
    }
    settings_path = project_path / "narration-settings.json"
    settings_path.write_text(json.dumps(settings, ensure_ascii=False, indent=2), encoding="utf-8")

    provider = settings["provider"]
    voice_id = settings["voiceId"] or "<voice_id>"
    rate = settings["rate"]
    project_arg = shlex.quote(str(project_path))
    voice_arg = shlex.quote(str(voice_id))
    rate_arg = shlex.quote(str(rate))
    if provider == "edge":
        command = f"python3 scripts/notes_to_audio.py {project_arg} --voice {voice_arg} --rate {rate_arg}"
    else:
        command = f"python3 scripts/notes_to_audio.py {project_arg} --provider {provider} --voice-id {voice_arg}"
    readme_path = audio_dir / "README.md"
    readme_path.write_text(
        (
            "# Narration handoff\n\n"
            f"- Provider: `{provider}`\n"
            f"- Voice: `{voice_id}`\n"
            f"- Rate: `{rate}`\n"
            f"- Notes: `{notes_dir}`\n\n"
            "Generate page-level narration audio from the repository root:\n\n"
            f"```bash\n{command}\n```\n"
        ),
        encoding="utf-8",
    )
    return [settings_path, readme_path]


def formal_quality_gate() -> dict[str, Any]:
    return {
        "level": "formal-business",
        "requiredInputs": [
            "visual direction pack or documented custom benchmark",
            "brand assets or explicit fallback strategy",
            "traceable source evidence",
            "page role and recipe contract with visual weight, layout family, page recipe, visual layer, raster policy, asset requirement, and anti-patterns",
            "ChatGPT-generation-first visual asset plan with prompts, filenames, and insertion targets",
            "small reusable element kit plan for section dividers, metric badges, process nodes, connectors, and icons",
            "page visual layer plan for no-text backgrounds, patterns, process accents, and metric accents",
            "public asset search plan for evidence/official references or explicit no-search rationale",
            "image/chart plan or explicit no-image strategy",
            "page rhythm and infographic strategy",
        ],
        "acceptanceCriteria": [
            "do not build the whole deck from headings and repeated cards only",
            "each page has a declared narrative role, layout family, and page recipe before generation",
            "three consecutive non-anchor pages must not share the same layout family or page recipe unless intentionally documented",
            "formal PPTX body pages remain editable and must not be replaced by full-page generated images",
            "ChatGPT/OpenAI is treated as the primary visual asset engine for custom supporting visuals",
            "small generated micro-assets are planned in visual-element-kit.md and reused across the deck",
            "page visual layers are no-text support assets and are not allowed to contain policy wording, numbers, logos, QR codes, or body copy",
            "generated assets are stored under assets/generated and listed in asset-plan.md",
            "public web asset searches are limited to evidence, official references, or brand boundaries and record source URLs, licensing notes, and insertion targets",
            "PPTX keeps real editable text, shapes, charts, and notes",
            "Web Deck has a complete visual system, layout variety, and desktop/mobile readability",
            "logo must not degrade into text fragments",
            "run formal delivery audit before delivery",
        ],
        "artifactChecks": [
            "manifest.json contains formal-business qualityGate",
            "design_spec.md and spec_lock.md contain visual direction, page role, recipe, visual layer, and raster policy contracts",
            "HTML/PPTX expose enough layout types",
            "real image/brand assets are used or no-image strategy is explicit",
            "asset-plan.md records public searches, generated assets, citations, and insert targets",
            "visual-element-kit.md records the reusable ChatGPT-generated micro-assets",
            "assets/generated/page-visuals/manifest.json records no-text page visual layers",
            "PPTX contains editable text objects",
            "no b/c-style logo text fragments",
        ],
        "reviewCommands": [
            "python3 scripts/audit_formal_delivery.py <project_path>",
            "python3 scripts/audit_design_completion.py <project_path>",
            "python3 scripts/audit_visual_recipes.py <project_path>",
        ],
        "assetStrategy": {
            "mode": "chatgpt-generation-first",
            "visualStrategyMode": "hybrid-editable",
            "rasterSlideMode": "disabled_for_formal_body",
            "brand": "Desktop draft uses a clean text fallback plus style palette; replace with approved brand assets for final delivery.",
            "primaryEngine": "For production, Codex should use ChatGPT/OpenAI image generation as the default visual asset engine for custom scenes and reusable micro-assets.",
            "microAssets": "Generate visual-element-kit.md assets before final assembly: section dividers, metric badges, process nodes, connectors, icon accents, textures, and callouts.",
            "pageVisualLayers": "Generate no-text page visual layers for backgrounds, patterns, process accents, and metric accents; keep text, numbers, charts, tables, and official assets editable.",
            "publicSearch": "Use public web sources mainly for evidence, official references, and brand boundaries; record source URL, license note, and insertion target.",
            "generatedAssets": "Store ChatGPT/OpenAI outputs under assets/generated/ and record prompts in asset-plan.md.",
            "images": "Explicit no-image strategy for this local draft: use editable shapes, charts, and information architecture unless source/approved or generated images are supplied.",
        },
    }


def formal_quality_profile(mode: str) -> dict[str, Any]:
    expected = [
        "sources/source.md",
        "design_spec.md",
        "spec_lock.md",
        "manifest.json",
        "project-brief.json",
        "quality-report.json",
        "design-quality-report.md",
        "assets/generated/page-visuals/manifest.json",
        "images/page_visual_prompts.md",
    ]
    expected.append("outputs/ultimate-ppt-master-preview.pptx" if mode == "pptx" else "outputs/index.html")
    return {
        "label": "正式商务交付质量",
        "userLevel": "中文办公用户 / formal business users",
        "acceptanceCriteria": [
            "核心结论清楚，页面节奏有变化",
            "证据和资料边界可追溯",
            "正式交付前通过 formal-business audit",
            "输出文件保留可编辑结构或明确 Web Deck 视觉策略",
        ],
        "expectedArtifacts": expected,
        "reviewCommands": [
            "python3 scripts/audit_formal_delivery.py <project_path>",
            "python3 scripts/audit_design_completion.py <project_path>",
            "python3 scripts/audit_visual_recipes.py <project_path>",
        ],
        "notFor": [
            "只需要快速标题草稿的场景",
            "需要直接发布真实品牌资产但尚未提供授权素材的场景",
        ],
    }


def codex_task_text(
    title: str,
    output_mode: str,
    quality_gate: dict[str, Any],
    workflow_state: dict[str, str],
    expected_artifacts: list[str],
) -> str:
    gate_inputs = "\n".join(f"- {item}" for item in quality_gate["requiredInputs"])
    gate_criteria = "\n".join(f"- {item}" for item in quality_gate["acceptanceCriteria"])
    gate_checks = "\n".join(f"- {item}" for item in quality_gate["artifactChecks"])
    artifacts = "\n".join(f"- {item}" for item in expected_artifacts)
    commands = "\n".join(f"- {item}" for item in quality_gate["reviewCommands"])
    return f"""# Codex Task

Project: {title}
Output mode: {output_mode}
Current workflow step: {workflow_state["currentStep"]}
Blocked reason: {workflow_state["blockedReason"] or "none"}

## Read First
1. AGENTS.md
2. manifest.json
3. project-brief.json
4. design_spec.md
5. spec_lock.md
6. quality-checklist.md if present
7. asset-plan.md
8. visual-element-kit.md
9. images/page_visual_prompts.md
10. sources/source.md

## Formal Business Gate
Required inputs:
{gate_inputs}

Acceptance criteria:
{gate_criteria}

Artifact checks:
{gate_checks}

## Asset Workflow
1. Inspect sources/source.md and generated previews before searching.
2. Treat ChatGPT/OpenAI as the primary visual asset engine: generate custom slide visuals and small reusable elements before final assembly.
3. From the Ultimate PPT Master repository root, run: `python3 scripts/generate_visual_element_kit.py <project_path>`.
4. From the Ultimate PPT Master repository root, run: `python3 scripts/generate_visual_layers.py <project_path>`.
5. If no IMAGE_BACKEND/OpenAI key is configured, do not block: use the Needs-Manual prompts in images/image_prompts.md and images/page_visual_prompts.md with ChatGPT and save outputs to the listed outputPath.
6. Generate the visual-element-kit.md micro-assets: section divider, metric badge, process node, connector, icon accent, subtle texture, and callout sticker.
7. Generate only no-text page visual layers: backgrounds, patterns, process accents, metric accents, or schematic atmosphere.
8. Save generated assets under assets/generated/ and insert them as real image objects, not full-slide screenshots.
9. Use public web search mainly for factual evidence, official references, brand boundaries, or source citations.
10. Record every generated prompt and every public source/license note in asset-plan.md.
11. Keep PPTX text, charts, tables, and labels editable wherever possible.

## Production Steps
1. Lock visual direction, brand/fallback strategy, evidence boundaries, page roles, page recipes, visual layers, raster policy, page rhythm, layout families, infographic strategy, asset-plan.md, and visual-element-kit.md.
2. Produce the requested {output_mode.upper()} delivery using the Ultimate PPT Master Skill workflow.
3. Vary layouts across narrative, comparison, timeline/process, metric, decision, and closing pages.
4. Do not let logos degrade into b/c-style text fragments.
5. Do not use full-page generated images for formal PPTX body pages unless explicitly approved and recorded in raster_policy.
6. Run review commands and update quality-report.json plus design-quality-report.md before final response.

## Expected Artifacts
{artifacts}

## Review Commands
{commands}

Final response: list generated files, ChatGPT micro-assets inserted, public references used, checks run, and remaining risks.
"""


def codex_agent_guide_text(quality_gate: dict[str, Any]) -> str:
    return f"""# AGENTS.md

## Codex Local Rules
- Work in this desktop project and the Ultimate PPT Master repository scripts only.
- Read codex-task.md before editing or generating deliverables.
- Keep private source material, customer data, internal screenshots, and API keys local unless the user explicitly approves upload.
- ChatGPT/OpenAI image generation is the primary visual asset engine. Read visual-element-kit.md and run or handle scripts/generate_visual_element_kit.py first when the deck needs visual richness.
- Use scripts/generate_visual_layers.py for no-text page visual layers. Body-page text, numbers, charts, tables, QR codes, and logos must stay editable or sourced.
- If no image backend/key is configured, use images/image_prompts.md and images/page_visual_prompts.md Needs-Manual prompts with ChatGPT and save outputs under assets/generated/.
- Public web search is allowed mainly for evidence, official references, and brand boundaries; record sources, usage notes, and insertion targets in asset-plan.md.
- Save generated outputs under assets/generated/ and record prompts in asset-plan.md.
- For level {quality_gate["level"]}, do not finish with repeated title-card slides, flat PPTX screenshots, broken logo fragments, or missing quality-report.json.
- Run the formal delivery audit, design completion audit, and visual recipe audit before final response whenever HTML or PPTX artifacts exist.
"""


def asset_plan_text(title: str, source_name: str, quality_gate: dict[str, Any]) -> str:
    gate_inputs = "\n".join(f"- {item}" for item in quality_gate["requiredInputs"])
    return f"""# Asset Plan

Project: {title}
Source: {source_name}

## Required Inputs
{gate_inputs}

## ChatGPT Generated Assets / ChatGPT 生成素材
- [ ] Treat ChatGPT/OpenAI as the primary visual asset engine: generate custom slide visuals and small reusable elements first, then use public search for evidence or official references.
- [ ] Generate the visual-element-kit.md micro-assets: section dividers, metric badges, process nodes, connectors, icon accents, subtle textures, and callout stickers.
- [ ] Store generated bitmap assets under assets/generated/ and generated SVG/icons under assets/generated/svg/.
- [ ] Record prompt, filename, target slide, and any manual edits before insertion.

## Public Asset Search / 公开素材检索
- [ ] Inspect sources/source.md and existing outputs before searching.
- [ ] Use public web search mainly for factual evidence, official references, brand boundaries, or source citations.
- [ ] Record source URL, publisher, license/usage note, and intended slide or Web Deck section.
- [ ] Do not upload private source files, customer data, or internal screenshots unless the user explicitly approves it.

## Insertion Targets
| Slide / Section | Need | Source or prompt | File path | Status |
| --- | --- | --- | --- | --- |
| Cover | brand-safe hero visual or approved logo treatment | pending | pending | pending |
| Evidence page | image/screenshot/chart supporting a claim | pending | pending | pending |
| Process page | editable diagram or generated scene | pending | pending | pending |

## Delivery Rule
Every inserted image, icon, logo, screenshot, or generated visual must be listed here before final delivery. If no imagery is used, keep the explicit no-image strategy in manifest.json / project-brief.json.
"""


def visual_element_kit_text(title: str, quality_gate: dict[str, Any]) -> str:
    mode = quality_gate.get("assetStrategy", {}).get("mode", "chatgpt-generation-first")
    return f"""# Visual Element Kit

Project: {title}
Mode: {mode}

## Purpose
Use ChatGPT/OpenAI as the primary visual asset engine. Generate small reusable elements before final slide production so the deck has a coherent visual language without relying on random stock imagery.

## Micro-assets / small reusable elements
| Asset type | Quantity target | Use | Output path | Status |
| --- | ---: | --- | --- | --- |
| section divider | 3 | chapter breaks and transition slides | assets/generated/dividers/ | pending |
| metric badge | 6 | KPI callouts, rights/benefits numbers, scorecards | assets/generated/badges/ | pending |
| process node | 5 | flow pages, timelines, service journey diagrams | assets/generated/process/ | pending |
| connector | 6 | arrows, dotted links, handoff paths, causal chains | assets/generated/connectors/ | pending |
| icon accent | 8 | small semantic markers for evidence, risk, action, user, channel | assets/generated/icons/ | pending |
| subtle pattern or texture | 2 | cover, section backgrounds, low-contrast visual depth | assets/generated/patterns/ | pending |
| callout sticker | 4 | reminders, caveats, delivery notes, decision highlights | assets/generated/callouts/ | pending |

## Prompt Pattern
- "Create a clean business presentation metric badge for [theme], flat vector-like style, transparent background, no text, colors aligned to [palette]."
- "Create a small process node icon for [step], formal government/finance presentation style, isolated object, no text, editable-friendly shape language."
- "Create a subtle abstract background texture for [theme], low contrast, no letters, no logos, widescreen presentation use."

## Insertion Rules
- Use micro-assets to enrich the visual system; keep body copy, numbers, and chart labels editable.
- Insert generated assets into PPTX/Web Deck as real images or graphic objects, not full-slide screenshots.
- Use public search for evidence, official references, and brand boundaries; use ChatGPT-generated micro-assets for the visual language.
- Register every generated file in asset-plan.md with prompt, filename, target slide, and insertion status.
"""


def visual_element_prompt_fallback_text(title: str) -> str:
    return f"""# ChatGPT Visual Element Prompts

Project: {title}
Status: Needs-Manual

Run this from the Ultimate PPT Master repository root when Codex takes over:

```bash
python3 scripts/generate_visual_element_kit.py <project_path>
```

If no IMAGE_BACKEND/OpenAI key is configured, paste the prompts implied by visual-element-kit.md into ChatGPT and save outputs under assets/generated/ before final assembly.

Required element types: section-divider, metric-badge, process-node, connector, icon-accent, subtle-pattern, callout-sticker.
Do not embed body text, numbers, logos, or long labels inside generated images; keep those editable in PPTX/Web Deck.
"""


def visual_direction_for_style(style: str, output_mode: str) -> dict[str, str]:
    if style == "consulting":
        direction_id = "consulting_analysis"
        benchmark = "Top-tier consulting analysis deck with varied evidence, process, and decision pages."
    elif style == "academic":
        direction_id = "training_courseware"
        benchmark = "Structured courseware deck with clear learning rhythm and diagram-first explanation pages."
    elif style in {"editorial", "swiss"} or output_mode == "web":
        direction_id = "launch_promotion"
        benchmark = "High-rhythm web presentation with strong anchors, metric moments, and controlled visual contrast."
    else:
        direction_id = "finance_internal_report"
        benchmark = "Formal internal business report with source-grounded evidence pages and restrained brand expression."
    return {
        "id": direction_id,
        "benchmark": benchmark,
        "release_boundary": "desktop-draft; replace fallback brand assets before formal external release",
    }


def infer_page_role(index: int, total: int, title: str, bullets: list[Any]) -> str:
    if index == 0:
        return "anchor"
    text = f"{title} {' '.join(str(item) for item in bullets)}".lower()
    if index == total - 1:
        return "closing"
    if re.search(r"流程|路径|步骤|办理|申请|开通|推进|workflow|process|step|pipeline", text):
        return "process"
    if re.search(r"权益|数字|指标|数据|结果|kpi|metric|benefit|output", text):
        return "benefit"
    if re.search(r"风险|提醒|边界|问题|疑问|注意|risk|issue|question|caveat", text):
        return "risk"
    if re.search(r"对比|比较|不同|差异|comparison|compare|matrix", text):
        return "comparison"
    if re.search(r"计划|行动|下一步|落地|安排|next|action|roadmap", text):
        return "action"
    if index == 1:
        return "context"
    return "evidence"


def page_contract_for_outline(outline: list[dict[str, Any]]) -> list[dict[str, str]]:
    role_recipes = {
        "anchor": ("cover_brand", "cover_brand.hero_left_visual"),
        "context": ("statement_plus_evidence", "statement_plus_evidence.left_rule_panel"),
        "comparison": ("comparison_matrix", "comparison_matrix.two_column_delta"),
        "process": ("process_flow", "process_flow.horizontal_steps"),
        "benefit": ("metric_panel", "metric_panel.large_number_strip"),
        "risk": ("risk_callout", "risk_callout.qa_stack"),
        "action": ("action_roadmap", "action_roadmap.owner_timeline"),
        "closing": ("closing_commitment", "closing_commitment.brand_tail"),
        "evidence": ("evidence_board", "evidence_board.source_table"),
    }
    fallback_recipes = [
        ("statement_plus_evidence", "statement_plus_evidence.left_rule_panel"),
        ("comparison_matrix", "comparison_matrix.two_column_delta"),
        ("process_flow", "process_flow.horizontal_steps"),
        ("metric_panel", "metric_panel.large_number_strip"),
        ("evidence_board", "evidence_board.source_table"),
        ("action_roadmap", "action_roadmap.owner_timeline"),
        ("risk_callout", "risk_callout.qa_stack"),
    ]
    anti_by_layout = {
        "cover_brand": "fake-logo, generic-gradient-cover",
        "statement_plus_evidence": "plain-title-bullets, oversized-empty-card",
        "comparison_matrix": "unlabeled-columns, weak-contrast-table",
        "process_flow": "disconnected-icons, unclear-step-order",
        "metric_panel": "random-big-number, no-unit-metric",
        "decision_table": "dense-table-without-hierarchy",
        "action_roadmap": "timeline-without-owner-or-next-step",
        "risk_callout": "risk-hidden-in-footnote",
        "evidence_board": "2x2-card-grid, source-free-claim",
        "closing_commitment": "generic-thank-you-only",
    }
    layer_by_recipe = {
        "cover_brand.hero_left_visual": "generated-background | no-text | 16:9",
        "statement_plus_evidence.left_rule_panel": "subtle-pattern | no-text | 16:9",
        "comparison_matrix.two_column_delta": "none",
        "process_flow.horizontal_steps": "generated-process-accent | no-text | 16:9",
        "metric_panel.large_number_strip": "generated-metric-accent | no-text | 16:9",
        "risk_callout.qa_stack": "none",
        "action_roadmap.owner_timeline": "schematic | no-text | 16:9",
        "closing_commitment.brand_tail": "generated-background | no-text | 16:9",
        "evidence_board.source_table": "none",
    }
    contracts: list[dict[str, str]] = []
    recent_layouts: list[str] = []
    recent_recipes: list[str] = []
    total = len(outline)

    for index, slide in enumerate(outline):
        page = f"P{index + 1:02d}"
        role = infer_page_role(index, total, str(slide.get("title", "")), list(slide.get("bullets", [])))
        layout, recipe = role_recipes.get(role) or fallback_recipes[(index - 1) % len(fallback_recipes)]
        if index > 0 and len(recent_layouts) >= 2 and recent_layouts[-1] == layout and recent_layouts[-2] == layout:
            layout, recipe = fallback_recipes[index % len(fallback_recipes)]
        if index > 0 and len(recent_recipes) >= 2 and recent_recipes[-1] == recipe and recent_recipes[-2] == recipe:
            layout, recipe = fallback_recipes[(index + 1) % len(fallback_recipes)]
        recent_layouts.append(layout)
        recent_recipes.append(recipe)
        if role == "anchor":
            visual_weight = "hero"
            asset_requirement = "real-logo-or-text-fallback"
            raster_policy = "allowed-cover"
        elif role in {"process", "benefit", "comparison"}:
            visual_weight = "high"
            asset_requirement = "schematic"
            raster_policy = "prohibited-formal-body"
        elif role in {"risk", "action"}:
            visual_weight = "medium-high"
            asset_requirement = "none-or-schematic"
            raster_policy = "prohibited-formal-body"
        elif role == "closing":
            visual_weight = "medium"
            asset_requirement = "real-logo-or-text-fallback"
            raster_policy = "allowed-section-tail"
        else:
            visual_weight = "medium"
            asset_requirement = "none"
            raster_policy = "prohibited-formal-body"
        visual_layer = layer_by_recipe.get(recipe, "none")
        if visual_layer != "none":
            layer_name = visual_layer.split("|", 1)[0].strip()
            visual_layer = f"{visual_layer} | assets/generated/page-visuals/{page}-{layer_name}.png"
        contracts.append(
            {
                "page": page,
                "title": str(slide.get("title", "")),
                "page_role": role,
                "visual_weight": visual_weight,
                "layout_family": layout,
                "page_recipe_id": recipe,
                "asset_requirement": asset_requirement,
                "visual_layer": visual_layer,
                "raster_policy": raster_policy,
                "anti_patterns": anti_by_layout.get(layout, "repeated-card-grid"),
            }
        )
    return contracts


def design_spec_text(title: str, style: str, output_mode: str, direction: dict[str, str], contracts: list[dict[str, str]]) -> str:
    rows = "\n".join(
        "| {page} | {page_role} | {visual_weight} | {layout_family} | {page_recipe_id} | {asset_requirement} | {visual_layer} | {raster_policy} | {anti_patterns} |".format(**item)
        for item in contracts
    )
    outline_rows = "\n".join(
        "- {page}: {title} | role={page_role} | recipe={page_recipe_id} | layer={visual_layer} | raster={raster_policy}".format(**item)
        for item in contracts
    )
    return f"""# Design Specification

## I. Project Info
- Title: {title}
- Output Mode: {output_mode}
- Style Preset: {style}

## II. Canvas
- Aspect Ratio: 16:9
- Delivery: desktop draft handoff for full Ultimate PPT Master refinement

## III. Visual Theme
- Visual Direction: {direction["id"]}
- Benchmark Sentence: {direction["benchmark"]}
- Top Aesthetic Risks: repeated card grids, fake logo fragments, text-only pages, weak evidence hierarchy
- External Release Boundary: {direction["release_boundary"]}
- Visual Strategy Mode: hybrid-editable
- Raster Slide Mode: disabled_for_formal_body

## IV. Typography
- Body baseline: 18 px equivalent
- Title/body split: display title plus readable dense body text

## V. Layout Principles
### Page Role / Visual Weight Contract
| page | page_role | visual_weight | layout_family | page_recipe_id | asset_requirement | visual_layer | raster_policy | anti_patterns |
| --- | --- | --- | --- | --- | --- | --- | --- | --- |
{rows}

## VI. Icon System
- Use restrained line icons only when they clarify a process, risk, or action.

## VII. Visualization Plan
- Use process flows, comparison matrices, metric panels, and evidence boards based on the page contract.
- Page recipes come from templates/page-recipes/index.json and drive PPTX/HTML builders.

## VIII. Image Resource List
- Default draft strategy: explicit no-image strategy unless approved, official, or generated assets are supplied.
- Cover uses real-logo-or-text-fallback until approved brand assets are available.
- Generated page visual layers must be no-text support assets and must not replace editable content.

## IX. Content Outline
{outline_rows}

## X. Speaker Notes
- Add speaker notes during full Agent refinement for risk, process, and benefit pages.

## XI. Technical Constraints
- Keep PPTX text editable.
- Avoid full-slide screenshots.
- Run audit_formal_delivery.py and audit_design_completion.py before handoff.
"""


def spec_lock_text(direction: dict[str, str], contracts: list[dict[str, str]]) -> str:
    def section(name: str, key: str) -> str:
        lines = [f"## {name}"]
        for item in contracts:
            lines.append(f"- {item['page']}: {item[key]}")
        return "\n".join(lines)

    return "\n\n".join(
        [
            "\n".join(
                [
                    "## visual_direction",
                    f"- id: {direction['id']}",
                    f"- benchmark: {direction['benchmark']}",
                    f"- release_boundary: {direction['release_boundary']}",
                ]
            ),
            section("page_roles", "page_role"),
            section("visual_weight", "visual_weight"),
            section("layout_family", "layout_family"),
            section("page_recipes", "page_recipe_id"),
            section("visual_layers", "visual_layer"),
            section("raster_policy", "raster_policy"),
            section("asset_requirements", "asset_requirement"),
            section("anti_patterns", "anti_patterns"),
        ]
    ) + "\n"


def page_visual_prompt(title: str, contract: dict[str, str]) -> str:
    return (
        f"Create a no-text visual support layer for page {contract['page']} of a formal business presentation titled "
        f"'{title}'. Page role: {contract['page_role']}. Page recipe: {contract['page_recipe_id']}. "
        f"Layer type: {contract['visual_layer'].split('|', 1)[0].strip()}. Use a polished official finance/service "
        "visual language. The image must contain no readable text, no numbers, no letters, no brand logos, no QR codes, "
        "no UI screenshots, no policy wording, and no watermark. It should support editable PPTX text, charts, tables, "
        "and shapes rather than replace them."
    )


def page_visual_layer_records(title: str, contracts: list[dict[str, str]]) -> tuple[dict[str, Any], dict[str, Any], str]:
    items: list[dict[str, Any]] = []
    prompt_items: list[dict[str, Any]] = []
    markdown_lines = [
        "# ChatGPT Page Visual Layer Prompts",
        "",
        "Status: Needs-Manual",
        "",
        "Use these prompts only for no-text support layers. Keep business copy, numbers, charts, tables, QR codes, and logos editable or sourced.",
        "",
    ]
    for contract in contracts:
        layer = contract["visual_layer"]
        if not layer or layer == "none":
            continue
        parts = [part.strip() for part in layer.split("|")]
        layer_type = parts[0]
        aspect = next((part for part in parts if re.fullmatch(r"\d+:\d+", part)), "16:9")
        output_path = next((part for part in parts if part.startswith("assets/") or part.startswith("images/")), f"assets/generated/page-visuals/{contract['page']}-{layer_type}.png")
        prompt = page_visual_prompt(title, contract)
        items.append(
            {
                "page": contract["page"],
                "layerType": layer_type,
                "textPolicy": "no-text",
                "aspectRatio": aspect,
                "prompt": prompt,
                "outputPath": output_path,
                "targetUse": "hybrid-editable page support layer",
                "status": "Needs-Manual",
                "backend": "manual-chatgpt",
                "insertedTargets": [],
            }
        )
        prompt_items.append(
            {
                "page": contract["page"],
                "layer_type": layer_type,
                "aspect_ratio": aspect,
                "image_size": "1K",
                "prompt": prompt,
                "output_path": output_path,
                "status": "Needs-Manual",
            }
        )
        markdown_lines.extend(
            [
                f"## {contract['page']} - {layer_type}",
                "",
                f"- Output path: `{output_path}`",
                f"- Aspect ratio: `{aspect}`",
                "- Status: `Needs-Manual`",
                "",
                prompt,
                "",
            ]
        )
    return (
        {
            "version": "page-visual-layers-v1",
            "mode": "hybrid-editable",
            "items": items,
        },
        {
            "version": "page-visual-prompts-v1",
            "status": "Needs-Manual",
            "items": prompt_items,
        },
        "\n".join(markdown_lines),
    )


def write_formal_delivery_files(
    project_path: Path,
    job: dict[str, Any],
    source_name: str,
    outline: list[dict[str, Any]],
    generated_files: list[str],
    preview_html: str,
    now: str,
) -> list[Path]:
    quality_gate = formal_quality_gate()
    quality_profile = formal_quality_profile(job["outputMode"])
    workflow_state = {
        "currentStep": "review",
        "blockedReason": "",
    }
    title = outline[0]["title"] if outline else source_name
    visual_direction = visual_direction_for_style(job["stylePreset"], job["outputMode"])
    page_contracts = page_contract_for_outline(outline)
    page_visual_manifest, page_visual_prompts, page_visual_prompt_md = page_visual_layer_records(str(title), page_contracts)
    design_spec_path = project_path / "design_spec.md"
    spec_lock_path = project_path / "spec_lock.md"
    manifest_path = project_path / "manifest.json"
    brief_path = project_path / "project-brief.json"
    report_path = project_path / "quality-report.json"
    codex_task_path = project_path / "codex-task.md"
    codex_guide_path = project_path / "AGENTS.md"
    asset_plan_path = project_path / "asset-plan.md"
    visual_element_kit_path = project_path / "visual-element-kit.md"
    page_visual_manifest_path = project_path / "assets" / "generated" / "page-visuals" / "manifest.json"
    page_visual_prompt_json_path = project_path / "images" / "page_visual_prompts.json"
    page_visual_prompt_md_path = project_path / "images" / "page_visual_prompts.md"
    image_prompt_fallback_path = project_path / "images" / "image_prompts.md"
    formal_file_paths = [
        design_spec_path,
        spec_lock_path,
        manifest_path,
        brief_path,
        report_path,
        codex_task_path,
        codex_guide_path,
        asset_plan_path,
        visual_element_kit_path,
        page_visual_manifest_path,
        page_visual_prompt_json_path,
        page_visual_prompt_md_path,
        image_prompt_fallback_path,
    ]
    brief = {
        "version": "desktop-worker-formal-v1",
        "title": title,
        "sourceName": source_name,
        "outputMode": job["outputMode"],
        "stylePreset": job["stylePreset"],
        "qualityProfile": quality_profile,
        "qualityGate": quality_gate,
        "workflowState": workflow_state,
        "storyboard": outline,
        "visualDirection": visual_direction,
        "pageContract": page_contracts,
        "visualStrategy": {
            "mode": "hybrid-editable",
            "rasterSlideMode": "disabled_for_formal_body",
            "benchmark": "bocom_social_card_private",
            "pageVisualManifest": "assets/generated/page-visuals/manifest.json",
        },
        "sourceStrategy": {
            "source": "sources/source.md",
            "privacy": "local-first; source files stay in the project folder",
            "evidenceBoundary": "desktop draft reflects available local text; Agent refinement should verify all numbers and claims against source files",
            "assetPlan": "asset-plan.md records public references, ChatGPT generated assets, prompts, source/license notes, and insertion targets",
            "visualElementKit": "visual-element-kit.md defines the ChatGPT-generation-first micro-assets to reuse across pages",
            "manualPromptFallback": "images/image_prompts.md records the Needs-Manual fallback until Codex runs generate_visual_element_kit.py",
        },
    }
    manifest = {
        "version": "desktop-worker-formal-v1",
        "createdAt": now,
        "app": "Ultimate PPT Master Desktop Worker",
        "projectPath": str(project_path),
        "sourceName": source_name,
        "outputMode": job["outputMode"],
        "stylePreset": job["stylePreset"],
        "generatedFiles": [*generated_files, *(str(path) for path in formal_file_paths)],
        "qualityProfile": quality_profile,
        "qualityGate": quality_gate,
        "workflowState": workflow_state,
        "expectedArtifacts": quality_profile["expectedArtifacts"],
        "reviewCommands": quality_gate["reviewCommands"],
        "artifactChecks": quality_gate["artifactChecks"],
        "visualDirection": visual_direction,
        "pageContractSummary": {
            "pages": len(page_contracts),
            "layoutFamilies": sorted({item["layout_family"] for item in page_contracts}),
            "pageRecipes": sorted({item["page_recipe_id"] for item in page_contracts}),
            "roles": sorted({item["page_role"] for item in page_contracts}),
        },
        "visualStrategy": {
            "mode": "hybrid-editable",
            "rasterSlideMode": "disabled_for_formal_body",
            "pageVisualManifest": "assets/generated/page-visuals/manifest.json",
        },
        "preview": {
            "hasHtml": bool(preview_html),
            "layoutStrategy": "varied data-layout pages for Web Deck; varied editable slide structures for PPTX",
        },
    }
    report = {
        "version": "desktop-worker-formal-v1",
        "status": "pending",
        "createdAt": now,
        "qualityProfile": quality_profile,
        "qualityGate": quality_gate,
        "workflowState": workflow_state,
        "reviewCommands": quality_gate["reviewCommands"],
        "visualDirection": visual_direction,
        "pageContractSummary": {
            "pages": len(page_contracts),
            "layoutFamilies": sorted({item["layout_family"] for item in page_contracts}),
            "pageRecipes": sorted({item["page_recipe_id"] for item in page_contracts}),
            "roles": sorted({item["page_role"] for item in page_contracts}),
        },
        "summary": {
            "zh": "桌面生成器已写入正式商务门禁、页面角色合约和页面配方合约。请运行 reviewCommands 中的正式商务审计、设计完成度审计和视觉配方审计；Design Doctor 默认只报告问题。",
            "en": "Desktop worker wrote the formal-business gate, page role contract, and page recipe contract. Run reviewCommands for delivery, design-completion, and visual-recipe audits; Design Doctor remains report-only by default.",
        },
        "checks": [
            {
                "id": "formal-business-gate",
                "status": "pending",
                "summary": "等待 audit_formal_delivery.py 验证。",
            }
        ],
    }
    files = [
        (manifest_path, manifest),
        (brief_path, brief),
        (report_path, report),
    ]
    written: list[Path] = []
    for path, payload in files:
        path.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        written.append(path)
    design_spec_path.write_text(design_spec_text(title, job["stylePreset"], job["outputMode"], visual_direction, page_contracts), encoding="utf-8")
    spec_lock_path.write_text(spec_lock_text(visual_direction, page_contracts), encoding="utf-8")
    written.extend([design_spec_path, spec_lock_path])
    page_visual_manifest_path.parent.mkdir(parents=True, exist_ok=True)
    page_visual_manifest_path.write_text(json.dumps(page_visual_manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    page_visual_prompt_json_path.parent.mkdir(parents=True, exist_ok=True)
    page_visual_prompt_json_path.write_text(json.dumps(page_visual_prompts, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    page_visual_prompt_md_path.write_text(page_visual_prompt_md, encoding="utf-8")
    written.extend([page_visual_manifest_path, page_visual_prompt_json_path, page_visual_prompt_md_path])
    codex_task_path.write_text(
        codex_task_text(title, job["outputMode"], quality_gate, workflow_state, quality_profile["expectedArtifacts"]),
        encoding="utf-8",
    )
    codex_guide_path.write_text(codex_agent_guide_text(quality_gate), encoding="utf-8")
    asset_plan_path.write_text(asset_plan_text(title, source_name, quality_gate), encoding="utf-8")
    visual_element_kit_path.write_text(visual_element_kit_text(title, quality_gate), encoding="utf-8")
    image_prompt_fallback_path.parent.mkdir(parents=True, exist_ok=True)
    image_prompt_fallback_path.write_text(visual_element_prompt_fallback_text(title), encoding="utf-8")
    written.extend([codex_task_path, codex_guide_path, asset_plan_path, visual_element_kit_path, image_prompt_fallback_path])
    return written


def write_runbook(project_path: Path, job: dict[str, Any], source_name: str) -> Path:
    runbook = project_path / "README.md"
    provider_config = job["providerConfig"]
    narration_line = (
        f"- Narration: `{provider_config['voiceProvider']}` / `{provider_config.get('voiceId') or 'none'}`\n"
        if provider_config.get("narrationEnabled")
        else "- Narration: `disabled`\n"
    )
    runbook.write_text(
        (
            "# Ultimate PPT Master Desktop Project\n\n"
            f"- Source: `{source_name}`\n"
            f"- Output mode: `{job['outputMode']}`\n"
            f"- Style preset: `{job['stylePreset']}`\n\n"
            "## Provider choices\n\n"
            f"- Text model provider: `{provider_config['modelProvider']}`\n"
            f"- Text model ID: `{provider_config.get('textModelId') or 'environment default'}`\n"
            f"- Image provider: `{provider_config['imageProvider']}`\n"
            f"- Image model ID: `{provider_config.get('imageModelId') or 'environment default'}`\n"
            f"{narration_line}\n"
            "This desktop MVP created immediate preview artifacts. For production-quality decks, "
            "open this project with an agent and follow the root `SKILL.md` workflow using the "
            "`sources/source.md` file as the source material.\n\n"
            "## Formal business gate\n\n"
            "This project writes `manifest.json`, `project-brief.json`, and `quality-report.json` "
            "with `qualityGate.level = formal-business`. It also writes `codex-task.md`, "
            "`AGENTS.md`, `asset-plan.md`, and `visual-element-kit.md` so Codex can use "
            "ChatGPT-generation-first visual assets, run `scripts/generate_visual_element_kit.py`, "
            "fall back to Needs-Manual prompts when no image key is configured, insert small "
            "reusable elements, and record prompts/sources before delivery. "
            "Validate before delivery:\n\n"
            "```bash\n"
            "python3 scripts/audit_formal_delivery.py <project_path>\n"
            "```\n"
        ),
        encoding="utf-8",
    )
    return runbook


def manifest_to_recent(manifest: dict[str, Any], manifest_path: Path) -> dict[str, Any]:
    source_name = str(manifest.get("sourceName") or manifest_path.parent.name)
    output_mode = str(manifest.get("outputMode") or ("web" if manifest.get("previewHtml") else "pptx"))
    updated_at = str(manifest.get("updatedAt") or datetime.fromtimestamp(manifest_path.stat().st_mtime, timezone.utc).isoformat().replace("+00:00", "Z"))
    return {
        "name": source_name,
        "mode": output_mode,
        "path": str(manifest.get("projectPath") or manifest_path.parent),
        "status": str(manifest.get("status") or "draft"),
        "createdAt": str(manifest.get("createdAt") or updated_at),
        "updatedAt": updated_at,
        "generatedFiles": manifest.get("generatedFiles") or [],
        "thumbnail": manifest.get("thumbnailSvg") or manifest.get("previewSvg") or "",
        "logsPath": manifest.get("logsPath") or "",
    }


def list_recent_projects(repo_root: Path, project_dir: str | None = None) -> list[dict[str, Any]]:
    root = Path(project_dir).expanduser() if project_dir else default_project_root(repo_root)
    if not root.exists():
        return []
    items: list[dict[str, Any]] = []
    for manifest_path in root.glob("*/desktop-manifest.json"):
        try:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            items.append(manifest_to_recent(manifest, manifest_path))
        except Exception:
            continue
    items.sort(key=lambda item: item["updatedAt"], reverse=True)
    return items[:12]


def run_job(job: dict[str, Any], repo_root: Path) -> dict[str, Any]:
    valid = validate_job(job)
    env = inspect_environment(repo_root)
    project_path = create_project_dir(valid, repo_root)
    text, source_name, source_extraction = source_to_text(valid, project_path, repo_root)
    outline = build_outline(text, valid["stylePreset"], valid["outputMode"])
    preview_svg = write_preview_svg(outline, project_path, valid["stylePreset"])
    generated_files: list[str] = []
    preview_html = ""

    if valid["outputMode"] == "pptx":
        generated_files.append(str(generate_pptx(outline, project_path, valid["stylePreset"])))
        generated_files.append(str(project_path / "preview" / "cover.svg"))
    else:
        html_path, preview_html = generate_web_deck(outline, project_path, valid["stylePreset"], repo_root)
        generated_files.append(str(html_path))

    narration_files = write_narration_handoff(project_path, outline, valid["providerConfig"])
    generated_files.extend(str(path) for path in narration_files)

    runbook = write_runbook(project_path, valid, source_name)
    generated_files.append(str(runbook))

    log_path = project_path / "logs" / "desktop-worker.log"
    log_path.write_text(
        "\n".join([
            f"provider: model={valid['providerConfig']['modelProvider']} text_model={valid['providerConfig'].get('textModelId') or 'default'}",
            f"provider: image={valid['providerConfig']['imageProvider']} image_model={valid['providerConfig'].get('imageModelId') or 'default'}",
            f"narration: enabled={valid['providerConfig']['narrationEnabled']} provider={valid['providerConfig']['voiceProvider']} voice={valid['providerConfig'].get('voiceId') or 'none'}",
            *[f"{step.key}: {step.message}" for step in STEPS],
        ]) + "\n",
        encoding="utf-8",
    )

    now = utc_now()
    formal_files = write_formal_delivery_files(project_path, valid, source_name, outline, generated_files, preview_html, now)
    generated_files.extend(str(path) for path in formal_files)
    recommendations = [recommend_settings(valid["source"])]
    checks = build_project_checks(env, valid, source_name)
    next_actions = build_next_actions(project_path, log_path, generated_files, valid["outputMode"])
    manifest = {
        "status": "complete",
        "projectPath": str(project_path),
        "logsPath": str(log_path),
        "generatedFiles": generated_files,
        "steps": [step.__dict__ for step in STEPS],
        "outputMode": valid["outputMode"],
        "stylePreset": valid["stylePreset"],
        "createdAt": now,
        "updatedAt": now,
        "recommendations": recommendations,
        "checks": checks,
        "nextActions": next_actions,
        "thumbnailSvg": preview_svg,
        "previewSvg": preview_svg,
        "previewHtml": preview_html,
        "outline": outline,
        "sourceName": source_name,
        "sourceExtraction": source_extraction,
        "providerConfig": valid["providerConfig"],
    }
    (project_path / "desktop-manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description="Ultimate PPT Master desktop worker")
    parser.add_argument("command", choices=["inspect", "run", "validate-job", "recommend", "list-projects"])
    parser.add_argument("--repo-root")
    parser.add_argument("--stdin", action="store_true", help="Read job JSON from stdin")
    parser.add_argument("--job", help="Path to a job JSON file")
    parser.add_argument("--project-dir", help="Project root for list-projects")
    args = parser.parse_args()

    repo_root = repo_root_from_args(args.repo_root)
    try:
        if args.command == "inspect":
            result = inspect_environment(repo_root)
        elif args.command == "list-projects":
            result = list_recent_projects(repo_root, args.project_dir)
        else:
            if args.stdin:
                job = read_json_from_stdin()
            elif args.job:
                job = json.loads(Path(args.job).read_text(encoding="utf-8"))
            else:
                raise ValueError("Use --stdin or --job for job commands.")
            if args.command == "validate-job":
                result = validate_job(job)
            elif args.command == "recommend":
                source = job.get("source", job)
                if not isinstance(source, dict):
                    raise ValueError("recommend requires a source object.")
                result = recommend_settings(source)
            else:
                result = run_job(job, repo_root)
        print(json.dumps(result, ensure_ascii=False))
        return 0
    except Exception as exc:
        print(json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
