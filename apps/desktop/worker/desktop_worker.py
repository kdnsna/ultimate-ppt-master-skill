#!/usr/bin/env python3
"""Desktop worker for the Ultimate PPT Master app.

The worker is intentionally small: it creates a local project folder, normalizes
the user's source into a first-pass outline, and emits preview artifacts that the
desktop app can show immediately. The full agentic PPT workflow still lives in
SKILL.md and scripts/.
"""

from __future__ import annotations

import argparse
import html
import json
import os
import platform
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception:  # pragma: no cover - covered by environment checks
    Presentation = None  # type: ignore
    RGBColor = None  # type: ignore
    PP_ALIGN = None  # type: ignore
    Inches = None  # type: ignore
    Pt = None  # type: ignore


OUTPUT_MODES = {"pptx", "web"}
STYLE_PRESETS = {"business", "consulting", "academic", "editorial", "swiss"}
SOURCE_KINDS = {"file", "text", "url", "markdown"}
TEXT_SUFFIXES = {".md", ".markdown", ".txt", ".csv", ".tsv", ".json"}
DOCUMENT_SUFFIXES = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx"}
DEFAULT_PROJECT_ROOT = "projects/desktop"


@dataclass(frozen=True)
class Step:
    key: str
    label: str
    message: str
    progress: int


STEPS = [
    Step("source", "资料处理", "读取源资料并创建本地项目。", 14),
    Step("strategy", "策略整理", "提取标题、关键句和页面骨架。", 32),
    Step("design", "设计锁定", "应用所选输出模式和风格预设。", 50),
    Step("generate", "生成预览", "生成可预览的 PPTX 或网页演示文件。", 74),
    Step("verify", "校验", "检查输出路径、预览内容和日志。", 90),
    Step("export", "导出", "保存结果并写入项目清单。", 100),
]


def sanitize_name(value: str) -> str:
    value = value.strip() or "untitled"
    safe = "".join(ch if ch.isalnum() or ch in "-_." else "-" for ch in value)
    safe = re.sub(r"-{2,}", "-", safe).strip("-_.")
    return safe[:80] or "untitled"


def repo_root_from_args(value: str | None) -> Path:
    if value:
        return Path(value).expanduser().resolve()
    return Path(__file__).resolve().parents[3]


def read_json_from_stdin() -> dict[str, Any]:
    raw = sys.stdin.read()
    if not raw.strip():
        raise ValueError("No JSON job payload received on stdin.")
    value = json.loads(raw)
    if not isinstance(value, dict):
        raise ValueError("Job payload must be a JSON object.")
    return value


def validate_job(job: dict[str, Any]) -> dict[str, Any]:
    source = job.get("source")
    if not isinstance(source, dict):
        raise ValueError("job.source is required.")

    source_kind = str(source.get("kind", "")).lower()
    source_value = str(source.get("value", "")).strip()
    if source_kind not in SOURCE_KINDS:
        raise ValueError(f"Unsupported source kind: {source_kind}")
    if not source_value:
        raise ValueError("source.value is required.")

    output_mode = str(job.get("outputMode", job.get("output_mode", ""))).lower()
    if output_mode not in OUTPUT_MODES:
        raise ValueError(f"Unsupported output mode: {output_mode}")

    style_preset = str(job.get("stylePreset", job.get("style_preset", ""))).lower()
    if style_preset not in STYLE_PRESETS:
        raise ValueError(f"Unsupported style preset: {style_preset}")

    project_dir = job.get("projectDir", job.get("project_dir"))
    if project_dir is not None and not isinstance(project_dir, str):
        raise ValueError("projectDir must be a string when provided.")

    return {
        "source": {
            "kind": source_kind,
            "value": source_value,
            "name": str(source.get("name") or "").strip(),
        },
        "outputMode": output_mode,
        "stylePreset": style_preset,
        "projectDir": project_dir,
        "providerConfig": job.get("providerConfig", job.get("provider_config")) or {},
    }


def utc_now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def source_suffix(source: dict[str, Any]) -> str:
    name = source.get("name") or source.get("value") or ""
    if source.get("kind") == "url":
        return ".url"
    return Path(str(name)).suffix.lower()


def recommend_settings(source: dict[str, Any]) -> dict[str, Any]:
    kind = str(source.get("kind", "")).lower()
    value = str(source.get("value", ""))
    name = str(source.get("name", ""))
    suffix = source_suffix(source)
    text = f"{name}\n{value}".lower()

    if kind == "url":
        return {
            "outputMode": "web",
            "stylePreset": "editorial",
            "pageRange": "6-10",
            "reason": "URL 内容通常更适合先做成可浏览的网页演示，便于快速预览和分享。",
        }
    if suffix in {".pptx", ".docx", ".xlsx", ".xlsm", ".pdf"}:
        return {
            "outputMode": "pptx",
            "stylePreset": "business",
            "pageRange": "8-14",
            "reason": "正式资料优先生成真实可编辑 PPTX，便于交付后继续修改。",
        }
    if "swiss" in text or "工程" in text or "product" in text or "架构" in text:
        return {
            "outputMode": "web",
            "stylePreset": "swiss",
            "pageRange": "7-12",
            "reason": "产品、工程和系统表达适合 Swiss Style 的网格、秩序和信息密度。",
        }
    if "演讲" in text or "发布" in text or "分享" in text or "keynote" in text:
        return {
            "outputMode": "web",
            "stylePreset": "editorial",
            "pageRange": "8-12",
            "reason": "演讲和发布场景更需要视觉记忆点，网页演示更有现场感。",
        }
    if "咨询" in text or "方案" in text or "strategy" in text:
        return {
            "outputMode": "pptx",
            "stylePreset": "consulting",
            "pageRange": "10-16",
            "reason": "方案材料适合结论先行的咨询风，并保留 PPTX 可编辑交付。",
        }
    if "课程" in text or "培训" in text or "学术" in text or "论文" in text:
        return {
            "outputMode": "pptx",
            "stylePreset": "academic",
            "pageRange": "12-20",
            "reason": "培训和学术内容需要章节清楚、证据稳定，PPTX 更适合二次讲授。",
        }
    return {
        "outputMode": "pptx",
        "stylePreset": "business",
        "pageRange": "8-12",
        "reason": "默认走可编辑 PPTX，覆盖多数正式汇报和团队交付场景。",
    }


def command_exists(name: str) -> bool:
    return shutil.which(name) is not None


def cairo_available() -> bool:
    if command_exists("pkg-config"):
        result = subprocess.run(
            ["pkg-config", "--exists", "cairo"],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
            check=False,
        )
        if result.returncode == 0:
            return True
    return command_exists("cairo-trace")


def inspect_environment(repo_root: Path) -> dict[str, Any]:
    venv_python = repo_root / ".venv" / "bin" / "python"
    return {
        "repoRoot": str(repo_root),
        "platform": platform.platform(),
        "python": {
            "executable": sys.executable,
            "version": platform.python_version(),
            "bundledVenv": venv_python.exists(),
            "pythonPptx": Presentation is not None,
        },
        "node": {
            "available": command_exists("node"),
            "npm": command_exists("npm"),
            "pnpm": command_exists("pnpm"),
        },
        "optional": {
            "cairo": cairo_available(),
            "rust": command_exists("rustc") and command_exists("cargo"),
        },
        "providers": {
            "openai": bool(os.environ.get("OPENAI_API_KEY")),
            "gemini": bool(os.environ.get("GEMINI_API_KEY")),
            "qwen": bool(os.environ.get("QWEN_API_KEY") or os.environ.get("DASHSCOPE_API_KEY")),
            "pexels": bool(os.environ.get("PEXELS_API_KEY")),
            "pixabay": bool(os.environ.get("PIXABAY_API_KEY")),
        },
    }


def build_project_checks(env: dict[str, Any], job: dict[str, Any], source_name: str) -> list[dict[str, str]]:
    suffix = Path(source_name).suffix.lower()
    checks = [
        {
            "key": "local-first",
            "label": "本地处理",
            "status": "ok",
            "detail": "源文件和输出都保存在本地项目目录，不上传云端。",
        },
        {
            "key": "editable-output",
            "label": "真实可编辑",
            "status": "ok" if job["outputMode"] == "pptx" else "warning",
            "detail": "PPTX 路线输出真实文字和形状；Web Deck 路线优先展示体验。",
        },
        {
            "key": "provider",
            "label": "模型 Provider",
            "status": "ok" if any(env["providers"].values()) else "warning",
            "detail": "未检测到 provider key。本轮预览可运行，生产级生成需要配置模型。",
        },
        {
            "key": "document-parser",
            "label": "资料解析",
            "status": "warning" if suffix in DOCUMENT_SUFFIXES else "ok",
            "detail": "PDF/DOCX/XLSX/PPTX 已进入项目目录；生产级全文解析由 Agent 工作流接管。",
        },
    ]
    if not env["python"]["bundledVenv"]:
        checks.append(
            {
                "key": "venv",
                "label": "Python venv",
                "status": "missing",
                "detail": "未找到仓库 .venv；请先按 INSTALL.md 初始化依赖。",
            }
        )
    if not env["optional"]["rust"]:
        checks.append(
            {
                "key": "rust",
                "label": "Tauri 打包",
                "status": "warning",
                "detail": "未检测到 Rust/Cargo。Web 壳可用，原生 .app/.dmg 打包前需要安装 Rust。",
            }
        )
    return checks


def build_next_actions(project_path: Path, log_path: Path, generated_files: list[str], mode: str) -> list[dict[str, str]]:
    actions = [
        {
            "key": "open-result",
            "label": "打开结果文件",
            "detail": "先检查封面、结构和风格是否符合交付场景。",
            "path": generated_files[0] if generated_files else str(project_path),
        },
        {
            "key": "open-folder",
            "label": "打开项目文件夹",
            "detail": "查看 sources、outputs、preview、logs 和 manifest。",
            "path": str(project_path),
        },
        {
            "key": "agent-handoff",
            "label": "交给 Agent 深加工",
            "detail": "使用 sources/source.md 和 SKILL.md 进入生产级生成流程。",
            "path": str(project_path / "README.md"),
        },
        {
            "key": "open-log",
            "label": "查看生成日志",
            "detail": "排查依赖、路径和导出问题。",
            "path": str(log_path),
        },
    ]
    if mode == "web":
        actions[0]["detail"] = "在浏览器中检查横向翻页、视觉节奏和分享效果。"
    return actions


def default_project_root(repo_root: Path) -> Path:
    return repo_root / DEFAULT_PROJECT_ROOT


def project_name_from_job(job: dict[str, Any]) -> str:
    source = job["source"]
    name = source["name"]
    if name:
        return sanitize_name(Path(name).stem)
    value = source["value"]
    if source["kind"] == "url":
        parsed = urlparse(value)
        return sanitize_name(parsed.netloc or "web-source")
    if source["kind"] == "file":
        return sanitize_name(Path(value).stem)
    first_line = value.splitlines()[0] if value.splitlines() else "desk"
    return sanitize_name(first_line[:32])


def create_project_dir(job: dict[str, Any], repo_root: Path) -> Path:
    base = Path(job["projectDir"]).expanduser() if job.get("projectDir") else default_project_root(repo_root)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    project_name = project_name_from_job(job)
    project_path = base / f"{project_name}-{job['outputMode']}-{timestamp}"
    for child in ("sources", "outputs", "logs", "preview"):
        (project_path / child).mkdir(parents=True, exist_ok=True)
    return project_path


def source_to_text(job: dict[str, Any], project_path: Path) -> tuple[str, str]:
    source = job["source"]
    kind = source["kind"]
    value = source["value"]
    sources_dir = project_path / "sources"

    if kind in {"text", "markdown"}:
        filename = "source.md" if kind == "markdown" else "source.txt"
        (sources_dir / filename).write_text(value, encoding="utf-8")
        return value, filename

    if kind == "url":
        (sources_dir / "source.url").write_text(value + "\n", encoding="utf-8")
        text = f"URL source: {value}\n\nUse the full agent workflow to fetch and ground this page before production export."
        (sources_dir / "source.md").write_text(text, encoding="utf-8")
        return text, "source.url"

    source_path = Path(value).expanduser()
    if not source_path.exists():
        raise FileNotFoundError(f"Source file not found: {source_path}")
    copied_path = sources_dir / source_path.name
    if source_path.resolve() != copied_path.resolve():
        shutil.copy2(source_path, copied_path)

    if source_path.suffix.lower() in TEXT_SUFFIXES:
        text = source_path.read_text(encoding="utf-8", errors="replace")
    else:
        text = (
            f"Imported file: {source_path.name}\n\n"
            "This desktop MVP created a local project shell. For full document extraction, "
            "run the full ultimate-ppt-master workflow from SKILL.md."
        )
    (sources_dir / "source.md").write_text(text, encoding="utf-8")
    return text, source_path.name


def looks_chinese(text: str) -> bool:
    return bool(re.search(r"[\u4e00-\u9fff]", text))


def extract_lines(text: str) -> list[str]:
    cleaned: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        line = re.sub(r"^#{1,6}\s*", "", line)
        line = re.sub(r"^[-*+]\s*", "", line)
        line = re.sub(r"^\d+[.)]\s*", "", line)
        if line and line not in cleaned:
            cleaned.append(line)
    if not cleaned:
        cleaned = ["Untitled presentation", "Add source material to generate a richer deck."]
    return cleaned[:18]


def build_outline(text: str, style: str, mode: str) -> list[dict[str, Any]]:
    lines = extract_lines(text)
    chinese = looks_chinese(text)
    title = lines[0][:80]
    if chinese:
        headings = {
            "context": "核心信息",
            "workflow": "生成路径",
            "output": "交付结果",
            "next": "下一步",
        }
        mode_label = "可编辑 PPTX" if mode == "pptx" else "杂志风网页 PPT"
    else:
        headings = {
            "context": "Core Signal",
            "workflow": "Generation Path",
            "output": "Delivery Output",
            "next": "Next Step",
        }
        mode_label = "Editable PPTX" if mode == "pptx" else "Magazine Web Deck"

    body = lines[1:] or lines
    return [
        {"title": title, "eyebrow": "Ultimate PPT Master", "bullets": [mode_label, f"Style preset: {style}"]},
        {"title": headings["context"], "eyebrow": "01", "bullets": body[:4]},
        {"title": headings["workflow"], "eyebrow": "02", "bullets": ["Import source", "Lock structure and style", "Generate preview", "Verify and export"]},
        {"title": headings["output"], "eyebrow": "03", "bullets": [mode_label, "Local project folder", "Preview artifact", "Runbook for full agent workflow"]},
        {"title": headings["next"], "eyebrow": "04", "bullets": ["Open the generated files", "Review page rhythm", "Continue with the full skill for production quality"]},
    ]


def write_preview_svg(outline: list[dict[str, Any]], project_path: Path, style: str) -> str:
    colors = {
        "business": ("#FFF8EC", "#172033", "#F97316"),
        "consulting": ("#F7F5EF", "#111827", "#2563EB"),
        "academic": ("#F4F8FF", "#172033", "#10B981"),
        "editorial": ("#FFF8EC", "#2F241B", "#C2410C"),
        "swiss": ("#F7FAFC", "#111827", "#E11D48"),
    }
    bg, ink, accent = colors.get(style, colors["business"])
    slide = outline[0]
    bullet_text = "".join(
        f'<text x="88" y="{430 + idx * 48}" fill="{ink}" font-family="Avenir Next, PingFang SC, sans-serif" font-size="30">{html.escape(str(item))}</text>'
        for idx, item in enumerate(slide["bullets"][:4])
    )
    svg = f"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1600 900">
  <rect width="1600" height="900" fill="{bg}"/>
  <path d="M-40 190C260 40 480 180 720 120C980 56 1120 194 1640 104" fill="none" stroke="{accent}" stroke-width="44" opacity="0.16"/>
  <rect x="70" y="70" width="1460" height="760" rx="44" fill="#FFFFFF" opacity="0.82"/>
  <text x="88" y="154" fill="{accent}" font-family="Avenir Next, PingFang SC, sans-serif" font-size="28" font-weight="800">{html.escape(slide['eyebrow'])}</text>
  <text x="88" y="292" fill="{ink}" font-family="Avenir Next, PingFang SC, sans-serif" font-size="72" font-weight="850">{html.escape(slide['title'])}</text>
  {bullet_text}
  <text x="88" y="750" fill="#667085" font-family="Avenir Next, PingFang SC, sans-serif" font-size="24">Desktop MVP preview · Full production workflow remains agent-driven</text>
</svg>"""
    path = project_path / "preview" / "cover.svg"
    path.write_text(svg, encoding="utf-8")
    return svg


def add_text_box(slide: Any, left: float, top: float, width: float, height: float, text: str, size: int, color: tuple[int, int, int], bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)


def generate_pptx(outline: list[dict[str, Any]], project_path: Path, style: str) -> Path:
    if Presentation is None:
        raise RuntimeError("python-pptx is not available. Install requirements.txt first.")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent_map = {
        "business": (249, 115, 22),
        "consulting": (37, 99, 235),
        "academic": (16, 185, 129),
        "editorial": (194, 65, 12),
        "swiss": (225, 29, 72),
    }
    accent = accent_map.get(style, accent_map["business"])
    for idx, item in enumerate(outline):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(255, 248, 236)
        add_text_box(slide, 0.62, 0.48, 2.2, 0.3, item["eyebrow"], 14, accent, True)
        add_text_box(slide, 0.62, 1.15, 11.3, 1.1, item["title"], 34 if idx == 0 else 30, (23, 32, 51), True)
        for bullet_idx, bullet in enumerate(item["bullets"][:5]):
            add_text_box(slide, 0.92, 2.7 + bullet_idx * 0.65, 10.8, 0.42, f"- {bullet}", 18, (73, 84, 104))
        shape = slide.shapes.add_shape(1, Inches(10.9), Inches(5.9), Inches(1.45), Inches(0.14))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*accent)
        shape.line.fill.background()
    output = project_path / "outputs" / "ultimate-ppt-master-preview.pptx"
    prs.save(output)
    return output


def generate_web_deck(outline: list[dict[str, Any]], project_path: Path, style: str) -> tuple[Path, str]:
    accent = "#E11D48" if style == "swiss" else "#F97316"
    sections = []
    for idx, slide in enumerate(outline):
        bullets = "\n".join(f"<li>{html.escape(str(item))}</li>" for item in slide["bullets"][:5])
        sections.append(
            f"""<section class="slide {'hero' if idx == 0 else ''}">
  <p class="eyebrow">{html.escape(slide['eyebrow'])}</p>
  <h1>{html.escape(slide['title'])}</h1>
  <ul>{bullets}</ul>
</section>"""
        )
    deck = f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Ultimate PPT Master Desktop Preview</title>
  <style>
    :root {{ color-scheme: light; --accent: {accent}; --ink: #172033; --muted: #667085; }}
    body {{ margin: 0; font-family: Avenir Next, PingFang SC, Microsoft YaHei, sans-serif; background: #fff8ec; color: var(--ink); }}
    main {{ display: flex; overflow-x: auto; scroll-snap-type: x mandatory; height: 100vh; }}
    .slide {{ min-width: 100vw; box-sizing: border-box; padding: 8vh 9vw; scroll-snap-align: start; display: grid; align-content: center; background: radial-gradient(circle at 18% 12%, rgba(249,115,22,.16), transparent 32%), linear-gradient(135deg, #fff8ec, #f4f8ff 58%, #fff3e8); }}
    .slide:nth-child(even) {{ background: radial-gradient(circle at 80% 18%, rgba(37,99,235,.16), transparent 30%), linear-gradient(135deg, #f8fbff, #fff8ec); }}
    .eyebrow {{ color: var(--accent); font-weight: 800; letter-spacing: .08em; text-transform: uppercase; }}
    h1 {{ font-size: clamp(44px, 7vw, 86px); line-height: .98; max-width: 980px; margin: 0 0 36px; }}
    ul {{ list-style: none; padding: 0; margin: 0; display: grid; gap: 16px; max-width: 820px; }}
    li {{ font-size: clamp(20px, 2.4vw, 30px); color: var(--muted); }}
    li::before {{ content: ""; display: inline-block; width: 14px; height: 14px; border-radius: 99px; background: var(--accent); margin-right: 14px; }}
  </style>
</head>
<body>
  <main>
    {' '.join(sections)}
  </main>
</body>
</html>"""
    output = project_path / "outputs" / "index.html"
    output.write_text(deck, encoding="utf-8")
    return output, deck


def write_runbook(project_path: Path, job: dict[str, Any], source_name: str) -> Path:
    runbook = project_path / "README.md"
    runbook.write_text(
        (
            "# Ultimate PPT Master Desktop Project\n\n"
            f"- Source: `{source_name}`\n"
            f"- Output mode: `{job['outputMode']}`\n"
            f"- Style preset: `{job['stylePreset']}`\n\n"
            "This desktop MVP created immediate preview artifacts. For production-quality decks, "
            "open this project with an agent and follow the root `SKILL.md` workflow using the "
            "`sources/source.md` file as the source material.\n"
        ),
        encoding="utf-8",
    )
    return runbook


def manifest_to_recent(manifest: dict[str, Any], manifest_path: Path) -> dict[str, Any]:
    source_name = str(manifest.get("sourceName") or manifest_path.parent.name)
    output_mode = str(manifest.get("outputMode") or ("web" if manifest.get("previewHtml") else "pptx"))
    updated_at = str(manifest.get("updatedAt") or datetime.fromtimestamp(manifest_path.stat().st_mtime, timezone.utc).isoformat().replace("+00:00", "Z"))
    return {
        "name": source_name,
        "mode": output_mode,
        "path": str(manifest.get("projectPath") or manifest_path.parent),
        "status": str(manifest.get("status") or "draft"),
        "createdAt": str(manifest.get("createdAt") or updated_at),
        "updatedAt": updated_at,
        "generatedFiles": manifest.get("generatedFiles") or [],
        "thumbnail": manifest.get("thumbnailSvg") or manifest.get("previewSvg") or "",
        "logsPath": manifest.get("logsPath") or "",
    }


def list_recent_projects(repo_root: Path, project_dir: str | None = None) -> list[dict[str, Any]]:
    root = Path(project_dir).expanduser() if project_dir else default_project_root(repo_root)
    if not root.exists():
        return []
    items: list[dict[str, Any]] = []
    for manifest_path in root.glob("*/desktop-manifest.json"):
        try:
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            items.append(manifest_to_recent(manifest, manifest_path))
        except Exception:
            continue
    items.sort(key=lambda item: item["updatedAt"], reverse=True)
    return items[:12]


def run_job(job: dict[str, Any], repo_root: Path) -> dict[str, Any]:
    valid = validate_job(job)
    env = inspect_environment(repo_root)
    project_path = create_project_dir(valid, repo_root)
    text, source_name = source_to_text(valid, project_path)
    outline = build_outline(text, valid["stylePreset"], valid["outputMode"])
    preview_svg = write_preview_svg(outline, project_path, valid["stylePreset"])
    generated_files: list[str] = []
    preview_html = ""

    if valid["outputMode"] == "pptx":
        generated_files.append(str(generate_pptx(outline, project_path, valid["stylePreset"])))
        generated_files.append(str(project_path / "preview" / "cover.svg"))
    else:
        html_path, preview_html = generate_web_deck(outline, project_path, valid["stylePreset"])
        generated_files.append(str(html_path))

    runbook = write_runbook(project_path, valid, source_name)
    generated_files.append(str(runbook))

    log_path = project_path / "logs" / "desktop-worker.log"
    log_path.write_text(
        "\n".join(f"{step.key}: {step.message}" for step in STEPS) + "\n",
        encoding="utf-8",
    )

    now = utc_now()
    recommendations = [recommend_settings(valid["source"])]
    checks = build_project_checks(env, valid, source_name)
    next_actions = build_next_actions(project_path, log_path, generated_files, valid["outputMode"])
    manifest = {
        "status": "complete",
        "projectPath": str(project_path),
        "logsPath": str(log_path),
        "generatedFiles": generated_files,
        "steps": [step.__dict__ for step in STEPS],
        "outputMode": valid["outputMode"],
        "stylePreset": valid["stylePreset"],
        "createdAt": now,
        "updatedAt": now,
        "recommendations": recommendations,
        "checks": checks,
        "nextActions": next_actions,
        "thumbnailSvg": preview_svg,
        "previewSvg": preview_svg,
        "previewHtml": preview_html,
        "outline": outline,
        "sourceName": source_name,
    }
    (project_path / "desktop-manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser(description="Ultimate PPT Master desktop worker")
    parser.add_argument("command", choices=["inspect", "run", "validate-job", "recommend", "list-projects"])
    parser.add_argument("--repo-root")
    parser.add_argument("--stdin", action="store_true", help="Read job JSON from stdin")
    parser.add_argument("--job", help="Path to a job JSON file")
    parser.add_argument("--project-dir", help="Project root for list-projects")
    args = parser.parse_args()

    repo_root = repo_root_from_args(args.repo_root)
    try:
        if args.command == "inspect":
            result = inspect_environment(repo_root)
        elif args.command == "list-projects":
            result = list_recent_projects(repo_root, args.project_dir)
        else:
            if args.stdin:
                job = read_json_from_stdin()
            elif args.job:
                job = json.loads(Path(args.job).read_text(encoding="utf-8"))
            else:
                raise ValueError("Use --stdin or --job for job commands.")
            if args.command == "validate-job":
                result = validate_job(job)
            elif args.command == "recommend":
                source = job.get("source", job)
                if not isinstance(source, dict):
                    raise ValueError("recommend requires a source object.")
                result = recommend_settings(source)
            else:
                result = run_job(job, repo_root)
        print(json.dumps(result, ensure_ascii=False))
        return 0
    except Exception as exc:
        print(json.dumps({"status": "error", "error": str(exc)}, ensure_ascii=False), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
