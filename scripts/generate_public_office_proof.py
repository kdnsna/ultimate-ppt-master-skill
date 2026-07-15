#!/usr/bin/env python3
"""Generate the public, synthetic, editable Executive Business Review PPTX proof."""

from __future__ import annotations

import argparse
import os
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PUBLIC_DIR = ROOT / "apps/web/public/examples/executive-business-review-starter"
EXAMPLE_DIR = ROOT / "examples/executive-business-review-starter"
DEFAULT_OUTPUT = PUBLIC_DIR / "executive-business-review-editable.pptx"

FONT = "Noto Sans CJK SC"
PAPER = "F6F3ED"
INK = "171714"
BLUE = "1D4ED8"
ORANGE = "D9573B"
GREEN = "73866C"
WHITE = "FFFFFF"
MUTED = "6B6A63"
LINE = "D9D4C9"
PALE_BLUE = "E8EEF9"
PALE_ORANGE = "F7E8E2"
PALE_GREEN = "E8ECE6"

CHART_XML_RE = re.compile(r"ppt/charts/chart\d+\.xml$")
CHART_AXIS_VALUE_RE = re.compile(rb'(<c:(?:axId|crossAx)\b[^>]*\bval=")(-\d+)(")')
CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
UINT32_MAX = (1 << 32) - 1


def verify_chart_axis_ids(path: Path) -> None:
    """Reject invalid UInt32 chart axis IDs or broken cross-axis references."""
    with zipfile.ZipFile(path) as package:
        bad_member = package.testzip()
        if bad_member:
            raise ValueError(f"Corrupt PPTX member after chart normalization: {bad_member}")
        for name in package.namelist():
            if not CHART_XML_RE.fullmatch(name):
                continue
            root = ET.fromstring(package.read(name))
            axis_ids = [int(node.attrib["val"]) for node in root.iter(f"{{{CHART_NS}}}axId")]
            cross_ids = [int(node.attrib["val"]) for node in root.iter(f"{{{CHART_NS}}}crossAx")]
            invalid = [value for value in (*axis_ids, *cross_ids) if not 0 <= value <= UINT32_MAX]
            if invalid:
                raise ValueError(f"{name}: chart axis ID is outside UInt32: {invalid}")
            missing = sorted(set(cross_ids) - set(axis_ids))
            if missing:
                raise ValueError(f"{name}: crossAx refers to unknown axis IDs: {missing}")


def normalize_chart_axis_ids(path: Path) -> int:
    """Normalize signed python-pptx axis IDs without rewriting the source in place."""
    replacements = 0
    temp_name: str | None = None
    try:
        with tempfile.NamedTemporaryFile(
            dir=path.parent,
            prefix=f".{path.stem}-",
            suffix=".pptx.tmp",
            delete=False,
        ) as temp_file:
            temp_name = temp_file.name

        with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(temp_name, "w", allowZip64=True) as target:
            target.comment = source.comment
            for info in source.infolist():
                data = source.read(info)
                if CHART_XML_RE.fullmatch(info.filename):

                    def replace_axis_value(match: re.Match[bytes]) -> bytes:
                        nonlocal replacements
                        signed_value = int(match.group(2))
                        replacements += 1
                        return match.group(1) + str(signed_value & UINT32_MAX).encode("ascii") + match.group(3)

                    data = CHART_AXIS_VALUE_RE.sub(replace_axis_value, data)
                target.writestr(info, data)

        temp_path = Path(temp_name)
        verify_chart_axis_ids(temp_path)
        os.chmod(temp_path, path.stat().st_mode)
        os.replace(temp_path, path)
        temp_name = None
        return replacements
    finally:
        if temp_name:
            Path(temp_name).unlink(missing_ok=True)


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def force_east_asian_font(font) -> None:
    """Set Latin and East Asian typefaces for PowerPoint and LibreOffice."""
    font.name = FONT
    properties = font._element
    for child in list(properties):
        if child.tag.endswith("}ea"):
            child.set("typeface", FONT)
            return
    east_asian = OxmlElement("a:ea")
    east_asian.set("typeface", FONT)
    properties.append(east_asian)


def set_fill(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_slide_background(slide, color: str) -> None:
    """Use the native slide background so canvas-fill shapes cannot overflow."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width: float = 1) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)


RADIUS_ADJUSTMENT = {
    "small": 0.04,
    "medium": 0.08,
    "large": 0.14,
    "pill": 0.5,
}


def add_rect(slide, x, y, w, h, fill=PAPER, line=None, radius=False):
    """Add a native surface using the deck's semantic corner-radius family."""
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    if line:
        set_line(shape, line)
    else:
        shape.line.fill.background()
    if radius:
        try:
            radius_token = "medium" if radius is True else str(radius)
            shape.adjustments[0] = RADIUS_ADJUSTMENT.get(radius_token, RADIUS_ADJUSTMENT["medium"])
        except (IndexError, ValueError):
            pass
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
    font: str = FONT,
    line_spacing: float = 1.08,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.font.name = font
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            force_east_asian_font(run.font)
    return shape


def add_rich_text(slide, runs, x, y, w, h, *, size=18, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for item in runs:
        run = paragraph.add_run()
        run.text = item[0]
        force_east_asian_font(run.font)
        run.font.size = Pt(item[1] if len(item) > 1 and item[1] else size)
        run.font.bold = bool(item[2]) if len(item) > 2 else False
        run.font.color.rgb = rgb(item[3] if len(item) > 3 else INK)
    return shape


def set_slide_title(slide, text: str, x, y, w, h, *, size=25, color=INK, bold=True):
    shape = slide.shapes.title
    if shape is None:
        return add_text(slide, text, x, y, w, h, size=size, color=color, bold=bold)
    shape.left, shape.top, shape.width, shape.height = Inches(x), Inches(y), Inches(w), Inches(h)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = rgb(color)
    for run in paragraph.runs:
        force_east_asian_font(run.font)
    tree = shape._element.getparent()
    tree.remove(shape._element)
    tree.append(shape._element)
    return shape


def add_chrome(slide, page: int, section: str, *, dark: bool = False) -> None:
    color = WHITE if dark else INK
    muted = "C9C7C0" if dark else MUTED
    add_text(slide, f"{page:02d}  /  {section}", 0.58, 0.28, 3.8, 0.24, size=8.5, color=muted, bold=True)
    add_text(slide, "ULTIMATE PPT MASTER · 公开合成案例", 8.85, 0.28, 3.9, 0.24, size=8.5, color=muted, bold=True, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.58), Inches(7.1), Inches(12.75), Inches(7.1))
    set_line(line, "4B4B47" if dark else LINE, 0.6)
    add_text(slide, "脱敏合成数据 · 仅用于可编辑 PPTX 质量验证", 0.58, 7.16, 7.5, 0.2, size=7.5, color=muted)
    add_text(slide, f"P{page:02d}", 11.8, 7.16, 0.95, 0.2, size=7.5, color=color, bold=True, align=PP_ALIGN.RIGHT)


def add_title(slide, title: str, takeaway: str, page: int, section: str) -> None:
    add_chrome(slide, page, section)
    set_slide_title(slide, title, 0.58, 0.67, 11.8, 0.5, size=25, bold=True)
    add_rect(slide, 0.58, 1.25, 0.08, 0.64, fill=BLUE)
    add_text(slide, takeaway, 0.82, 1.24, 11.1, 0.7, size=15, color=MUTED)


def add_note(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text
    for paragraph in notes.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(12)
        for run in paragraph.runs:
            force_east_asian_font(run.font)


def style_chart(chart, *, show_legend=True, max_scale=None, percent=False) -> None:
    chart.has_title = False
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        force_east_asian_font(chart.legend.font)
        chart.legend.font.size = Pt(9)
    chart.chart_style = 2
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = rgb(LINE)
    force_east_asian_font(chart.value_axis.tick_labels.font)
    chart.value_axis.tick_labels.font.size = Pt(9)
    force_east_asian_font(chart.category_axis.tick_labels.font)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.format.line.color.rgb = rgb(LINE)
    chart.value_axis.format.line.color.rgb = rgb(LINE)
    if max_scale is not None:
        chart.value_axis.maximum_scale = max_scale
    if percent:
        chart.value_axis.tick_labels.number_format = "0%"


def style_table(table, widths, header_fill=INK) -> None:
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for row_index, row in enumerate(table.rows):
        row.height = Inches(0.46 if row_index else 0.42)
        for cell in row.cells:
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(header_fill if row_index == 0 else (WHITE if row_index % 2 else PAPER))
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT
                paragraph.font.size = Pt(9.5 if row_index else 9)
                paragraph.font.bold = row_index == 0
                paragraph.font.color.rgb = rgb(WHITE if row_index == 0 else INK)
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    force_east_asian_font(run.font)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "季度经营复盘：从增长到可持续增长"
    prs.core_properties.subject = "脱敏合成数据的公开可编辑 PPTX 质量案例"
    prs.core_properties.author = ""
    prs.core_properties.last_modified_by = ""
    prs.core_properties.keywords = "editable PPTX, Executive Business Review, synthetic data"
    blank = prs.slide_layouts[5]

    # P01 — light-first report cover with one dominant soft evidence panel.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, PAPER)
    add_rect(slide, 0.58, 0.65, 0.1, 1.08, fill=BLUE, radius="small")
    add_text(slide, "EXECUTIVE BUSINESS REVIEW", 0.86, 0.68, 5.8, 0.28, size=9, color=BLUE, bold=True)
    set_slide_title(slide, "季度经营复盘", 0.86, 1.28, 6.6, 0.68, size=34, color=INK, bold=True)
    add_text(slide, "从增长，到可持续增长", 0.86, 2.08, 6.7, 0.6, size=25, color=MUTED)
    add_rect(slide, 0.86, 3.02, 0.82, 0.08, fill=ORANGE, radius="small")
    add_text(slide, "下季度的核心不是单纯扩大收入，\n而是同时恢复留存、保护毛利、提升管道质量。", 0.86, 3.38, 6.24, 1.2, size=16, color=INK, line_spacing=1.18)
    add_text(slide, "2026 Q2 · 公开脱敏合成版", 0.86, 5.92, 5.3, 0.3, size=9, color=MUTED)
    add_rect(slide, 7.95, 0.72, 4.55, 5.9, fill=WHITE, line=LINE, radius="large")
    add_text(slide, "本期经营信号", 8.42, 1.08, 2.4, 0.3, size=10, color=MUTED, bold=True)
    add_text(slide, "增长仍在，质量承压", 8.42, 1.48, 3.45, 0.45, size=19, color=INK, bold=True)
    metrics = [("12%", "收入增速", BLUE), ("82%", "客户留存", ORANGE), ("2.7×", "管道覆盖", GREEN)]
    for idx, (value, label, color) in enumerate(metrics):
        y = 2.25 + idx * 1.15
        add_text(slide, value, 8.42, y, 1.75, 0.52, size=27, color=color, bold=True)
        add_text(slide, label, 10.25, y + 0.09, 1.52, 0.28, size=10, color=MUTED, bold=True)
        if idx < 2:
            add_rect(slide, 8.42, y + 0.78, 3.55, 0.012, fill=LINE)
    add_rect(slide, 8.42, 5.86, 3.55, 0.38, fill=PALE_GREEN, radius="pill")
    add_text(slide, "下季度主线：先稳质量，再扩规模", 8.56, 5.97, 3.26, 0.17, size=8.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_chrome(slide, 1, "封面")
    add_note(slide, "这是一份完全由脱敏合成数据构成的公开案例。本期收入仍在增长，但留存和毛利同时承压，因此下季度需要把增长从量转向质。")

    # P02 — executive summary.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, PAPER)
    add_title(slide, "下季度必须同时修复留存、毛利和转化质量", "收入动能仍在，但三个结构性缺口已经影响增长的可持续性。", 2, "管理摘要")
    add_rect(slide, 0.58, 2.14, 12.17, 0.94, fill=PALE_BLUE, line="CBD8F2", radius="medium")
    add_text(slide, "首要判断", 0.88, 2.39, 1.25, 0.23, size=9, color=BLUE, bold=True)
    add_text(slide, "不先稳住留存，新客增量会被继续抵消。", 2.25, 2.32, 9.85, 0.36, size=18, color=INK, bold=True)
    cards = [
        ("01", "留存恢复", "30 天战情机制\n目标 ≥ 85%", ORANGE, PALE_ORANGE),
        ("02", "毛利保护", "调整渠道与折扣\n目标 ≥ 39%", BLUE, PALE_BLUE),
        ("03", "管道提纯", "商机分层与资格\n目标 ≥ 3.0×", GREEN, PALE_GREEN),
    ]
    for idx, (num, title, body, color, pale) in enumerate(cards):
        x = 0.58 + idx * 4.1
        add_rect(slide, x, 3.42, 3.84, 2.54, fill=WHITE, line=LINE, radius="medium")
        add_rect(slide, x + 0.08, 3.68, 0.08, 2.02, fill=color, radius="small")
        add_text(slide, num, x + 0.34, 3.75, 0.55, 0.32, size=11, color=color, bold=True)
        add_text(slide, title, x + 0.34, 4.18, 2.95, 0.36, size=18, bold=True)
        add_rect(slide, x + 0.34, 4.72, 0.62, 0.05, fill=color)
        add_text(slide, body, x + 0.34, 5.02, 3.05, 0.64, size=12.5, color=MUTED, line_spacing=1.18)
        add_rect(slide, x + 3.23, 3.68, 0.34, 0.34, fill=pale, radius="medium")
    add_note(slide, "管理摘要只保留三个决策项。第一，将留存恢复设为最高优先级。第二，用渠道和折扣约束保护毛利。第三，让管道从覆盖量转向高质量转化。")

    # P03 — KPI scorecard with a native table.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, PAPER)
    add_title(slide, "四项核心指标中，只有收入增速和管道覆盖较上期改善", "留存下降 2 个百分点，毛利下降 1 个百分点；这两项负向变化优先于继续扩张。", 3, "指标总览")
    metrics = [
        ("12%", "+3pp", "收入增速", "低于目标 3pp", BLUE, PALE_BLUE),
        ("38%", "-1pp", "毛利率", "渠道结构承压", ORANGE, PALE_ORANGE),
        ("82%", "-2pp", "客户留存", "首要风险", ORANGE, PALE_ORANGE),
        ("2.7×", "+0.3×", "管道覆盖", "转化质量分化", GREEN, PALE_GREEN),
    ]
    for idx, (value, delta, label, note, color, pale) in enumerate(metrics):
        x = 0.58 + idx * 3.06
        add_rect(slide, x, 2.15, 2.8, 1.78, fill=WHITE, line=LINE, radius=True)
        add_text(slide, label, x + 0.22, 2.42, 2.0, 0.24, size=9.5, color=MUTED, bold=True)
        add_text(slide, value, x + 0.22, 2.78, 1.35, 0.48, size=27, color=color, bold=True)
        add_rect(slide, x + 1.73, 2.83, 0.77, 0.31, fill=pale, radius=True)
        add_text(slide, delta, x + 1.73, 2.89, 0.77, 0.18, size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, note, x + 0.22, 3.47, 2.25, 0.23, size=8.5, color=MUTED)
    rows, cols = 5, 5
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.58), Inches(4.28), Inches(12.17), Inches(1.96))
    table = table_shape.table
    data = [
        ["指标", "本期", "上期", "目标", "管理判断"],
        ["收入增速", "12%", "9%", "15%", "增长改善，但挑战目标未达"],
        ["毛利率", "38%", "39%", "40%", "立即优化渠道结构"],
        ["客户留存", "82%", "84%", "88%", "进入 P0 战情机制"],
        ["管道覆盖", "2.7×", "2.4×", "3.0×", "扩容同时提高商机资格"],
    ]
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    style_table(table, [2.1, 1.15, 1.15, 1.15, 6.62])
    add_note(slide, "四项指标不能只看当前值。收入和管道比上期好，但仍未达目标；留存和毛利则是同比和目标双重缺口。因此资源排序应优先修复结构性风险。")

    # P04 — native clustered bar chart.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, WHITE)
    add_title(slide, "最大目标缺口来自留存，而不是需求不足", "留存距目标尚差 6 个百分点；管道覆盖已经接近目标，需要的是提纯而非单纯扩容。", 4, "目标差距")
    chart_data = ChartData()
    chart_data.categories = ["收入增速", "毛利率", "客户留存"]
    chart_data.add_series("本期", (12, 38, 82))
    chart_data.add_series("目标", (15, 40, 88))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.72), Inches(2.18), Inches(7.35), Inches(4.3), chart_data).chart
    style_chart(chart, show_legend=True, max_scale=100)
    chart.plots[0].gap_width = 78
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = rgb(BLUE)
    chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = rgb("C9C7C0")
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    force_east_asian_font(chart.plots[0].data_labels.font)
    chart.plots[0].data_labels.font.size = Pt(9)
    add_rect(slide, 8.53, 2.2, 3.72, 1.54, fill=PALE_ORANGE, line="EFC9BA", radius="large")
    add_text(slide, "6pp", 8.86, 2.5, 1.35, 0.44, size=26, color=ORANGE, bold=True)
    add_text(slide, "留存距目标的缺口", 10.1, 2.55, 1.75, 0.38, size=10, color=INK, bold=True)
    add_text(slide, "结论", 8.56, 4.18, 0.8, 0.25, size=9, color=BLUE, bold=True)
    add_text(slide, "需求仍然存在，\n执行质量决定增长能否留下。", 8.56, 4.54, 3.38, 0.82, size=16, bold=True, line_spacing=1.2)
    add_text(slide, "优先恢复留存；同时对管道做分层与资格检查。", 8.56, 5.55, 3.4, 0.65, size=11, color=MUTED)
    add_note(slide, "这张图把当前值和目标直接对比。留存的绝对缺口最大，距离目标六个百分点。管道覆盖已经接近目标，这说明当前不缺商机数量，缺的是商机质量和客户留存。")

    # P05 — driver chart.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, PAPER)
    add_title(slide, "新客获取拉动增长，但留存和渠道结构同时抵消成果", "本页是方向性归因，只展示脱敏合成资料中的已知边界，不延伸为完整财务拆分。", 5, "驱动分析")
    chart_data = ChartData()
    chart_data.categories = ["新客获取", "留存下降", "渠道结构", "管道扩容"]
    chart_data.add_series("方向性影响", (5, -2, -1, 0.3))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.72), Inches(2.18), Inches(8.1), Inches(4.15), chart_data).chart
    style_chart(chart, show_legend=False)
    chart.value_axis.minimum_scale = -3
    chart.value_axis.maximum_scale = 6
    chart.value_axis.major_unit = 1
    chart.plots[0].gap_width = 52
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = rgb(BLUE)
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    force_east_asian_font(chart.plots[0].data_labels.font)
    chart.plots[0].data_labels.font.size = Pt(9)
    add_rect(slide, 9.2, 2.18, 3.05, 0.1, fill=ORANGE)
    add_text(slide, "经营含义", 9.2, 2.56, 2.0, 0.28, size=10, color=ORANGE, bold=True)
    add_text(slide, "获客不是问题，\n客户能否留下才是。", 9.2, 3.04, 2.78, 0.85, size=18, bold=True, line_spacing=1.2)
    add_text(slide, "因此下季度的首项管理动作应从“追加获客”转向“留存战情 + 渠道治理”。", 9.2, 4.35, 2.86, 1.05, size=11.5, color=MUTED, line_spacing=1.22)
    add_rect(slide, 9.2, 5.73, 2.86, 0.46, fill=PALE_ORANGE, radius=True)
    add_text(slide, "证据边界：方向性合成归因", 9.35, 5.87, 2.52, 0.18, size=8, color=ORANGE, bold=True)
    add_note(slide, "这是方向性的驱动分析，不是完整的财务瀑布图。新客获取是主要正向驱动，留存和渠道结构则抵消了一部分成果。管理上的含义是，继续获客必须与留存和毛利治理同步。")

    # P06 — risk/opportunity matrix.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, WHITE)
    add_title(slide, "留存下降是高概率高影响风险，管道提纯是最快可转化的机会", "把管理注意力从“所有问题都重要”收缩为两个优先动作。", 6, "风险与机会")
    x0, y0, w, h = 0.95, 2.23, 8.15, 3.95
    add_rect(slide, x0, y0, w, h, fill=PAPER, line=LINE)
    v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0 + w / 2), Inches(y0), Inches(x0 + w / 2), Inches(y0 + h)); set_line(v, LINE, 1)
    hline = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0 + h / 2), Inches(x0 + w), Inches(y0 + h / 2)); set_line(hline, LINE, 1)
    add_text(slide, "高影响", 0.53, 2.15, 0.34, 1.0, size=8, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "高概率", 7.9, 6.33, 1.2, 0.22, size=8, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_rect(slide, 5.66, 2.69, 2.72, 0.86, fill=PALE_ORANGE, line=ORANGE, radius=True)
    add_text(slide, "留存持续下降", 5.9, 2.9, 2.2, 0.25, size=12, color=ORANGE, bold=True)
    add_rect(slide, 4.95, 4.45, 2.63, 0.86, fill=PALE_BLUE, line=BLUE, radius=True)
    add_text(slide, "渠道毛利承压", 5.18, 4.66, 2.12, 0.25, size=12, color=BLUE, bold=True)
    add_rect(slide, 1.7, 3.1, 2.68, 0.86, fill=PALE_GREEN, line=GREEN, radius=True)
    add_text(slide, "管道分层提纯", 1.94, 3.31, 2.18, 0.25, size=12, color=GREEN, bold=True)
    add_rect(slide, 9.55, 2.24, 2.65, 1.3, fill=PALE_ORANGE, line=ORANGE, radius="medium")
    add_text(slide, "P0", 9.85, 2.51, 0.48, 0.3, size=12, color=ORANGE, bold=True)
    add_text(slide, "留存战情", 10.37, 2.48, 1.45, 0.3, size=15, color=INK, bold=True)
    add_text(slide, "30 天内恢复至 ≥85%", 9.85, 3.0, 2.04, 0.25, size=9, color=MUTED)
    add_rect(slide, 9.55, 3.85, 2.65, 1.3, fill=PALE_GREEN, line=GREEN, radius=True)
    add_text(slide, "P1", 9.85, 4.12, 0.48, 0.3, size=12, color=GREEN, bold=True)
    add_text(slide, "管道提纯", 10.37, 4.09, 1.45, 0.3, size=15, bold=True)
    add_text(slide, "高质量覆盖达 3.0×", 9.85, 4.61, 2.04, 0.25, size=9, color=MUTED)
    add_text(slide, "中风险的毛利问题纳入同期渠道治理，不单独发起第三条并行战线。", 9.55, 5.55, 2.58, 0.7, size=9.5, color=MUTED)
    add_note(slide, "风险矩阵的目的是帮助管理层排序。留存是高概率且高影响的风险，需要立即进入战情机制。管道提纯的数据基础已经存在，是最快可转化的机会。毛利则应并入渠道治理，避免同时启动过多战线。")

    # P07 — priority action contracts.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, PAPER)
    add_title(slide, "三个行动全部绑定责任人、时限和可量化验收指标", "把经营建议改写为可追踪的交付合同，避免“下季度继续加强”。", 7, "优先行动")
    actions = [
        ("P0", "留存战情", "对高风险客户建立 30 天战情机制", "客户成功负责人", "第 2 周", "留存率 ≥85%", ORANGE),
        ("P0", "渠道治理", "重新设定渠道组合与折扣边界", "商务负责人", "第 4 周", "毛利率 ≥39%", BLUE),
        ("P1", "管道提纯", "将商机分层与资格检查纳入周会", "销售运营负责人", "第 3 周", "覆盖 ≥3.0×", GREEN),
    ]
    for idx, (priority, title, action, owner, deadline, metric, color) in enumerate(actions):
        y = 2.18 + idx * 1.36
        add_rect(slide, 0.58, y, 12.17, 1.08, fill=WHITE, line=LINE, radius=True)
        add_rect(slide, 0.66, y + 0.18, 0.08, 0.72, fill=color, radius=True)
        add_text(slide, priority, 0.93, y + 0.21, 0.58, 0.26, size=10, color=color, bold=True)
        add_text(slide, title, 1.66, y + 0.17, 1.62, 0.29, size=14, bold=True)
        add_text(slide, action, 3.44, y + 0.17, 3.76, 0.52, size=11.5, color=MUTED)
        add_text(slide, owner, 7.42, y + 0.17, 1.78, 0.5, size=10, bold=True)
        add_text(slide, deadline, 9.47, y + 0.17, 0.87, 0.28, size=9, color=MUTED)
        add_rect(slide, 10.47, y + 0.17, 1.92, 0.5, fill=(PALE_ORANGE if color == ORANGE else PALE_BLUE if color == BLUE else PALE_GREEN), radius=True)
        add_text(slide, metric, 10.57, y + 0.31, 1.72, 0.2, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "优先级", 0.93, 6.34, 0.85, 0.2, size=8, color=MUTED, bold=True)
    add_text(slide, "行动合同", 3.44, 6.34, 1.0, 0.2, size=8, color=MUTED, bold=True)
    add_text(slide, "责任人", 7.42, 6.34, 0.8, 0.2, size=8, color=MUTED, bold=True)
    add_text(slide, "时限", 9.47, 6.34, 0.55, 0.2, size=8, color=MUTED, bold=True)
    add_text(slide, "验收指标", 10.8, 6.34, 0.95, 0.2, size=8, color=MUTED, bold=True)
    add_note(slide, "这三个行动都不是抽象建议。每项都有明确责任人、最迟完成时间和验收指标。两个 P0 行动分别针对留存和毛利，P1 行动则把管道扩容转化为更高质量的销售机会。")

    # P08 — roadmap table.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, WHITE)
    add_title(
        slide,
        "前四周完成机制建立，之后用月度复核验证是否真正改善",
        "以下月度门槛为 Agent 合成假设，待业务确认；周节奏落地、月节奏复核。",
        8,
        "落地路径",
    )
    phases = [
        ("01", "第 1–2 周", "建立留存战情", "识别高风险客户\n明确挽回责任人", ORANGE),
        ("02", "第 3–4 周", "完成渠道与管道治理", "锁定折扣边界\n将商机分层纳入周会", BLUE),
        ("03", "第 2–3 月", "按月复核三项指标", "留存 · 毛利\n高质量管道覆盖", GREEN),
    ]
    for idx, (num, timing, title, body, color) in enumerate(phases):
        x = 0.72 + idx * 4.12
        add_text(slide, num, x, 2.27, 0.56, 0.3, size=10, color=color, bold=True)
        add_rect(slide, x, 2.72, 3.62, 0.1, fill=color)
        add_text(slide, timing, x, 3.07, 2.0, 0.27, size=10, color=MUTED, bold=True)
        add_text(slide, title, x, 3.49, 3.35, 0.55, size=16, bold=True)
        add_text(slide, body, x, 4.3, 3.18, 0.72, size=11, color=MUTED, line_spacing=1.2)
        if idx < 2:
            arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 3.66), Inches(2.77), Inches(x + 4.0), Inches(2.77)); set_line(arrow, LINE, 1.2)
    rows, cols = 4, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.72), Inches(5.48), Inches(11.98), Inches(1.25))
    table = table_shape.table
    table_data = [
        ["假设复核节点", "留存", "毛利", "高质量管道"],
        ["月末 1", "≥85%", "≥39%", "≥2.8×"],
        ["月末 2", "≥86%", "≥39.5%", "≥3.0×"],
        ["季末", "趋近 88%", "趋近 40%", "保持 ≥3.0×"],
    ]
    for r, row in enumerate(table_data):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
    style_table(table, [2.25, 3.2, 3.2, 3.33], header_fill=INK)
    for row in table.rows:
        row.height = Inches(0.31)
    add_note(
        slide,
        "落地路径分为三段。前两周建立留存战情，第三到第四周完成渠道与管道治理。"
        "表中月末和季末数值是为公开合成案例设定的 Agent 管理门槛假设，不是企业历史事实或已批准目标；"
        "正式使用前必须由业务负责人确认或替换。",
    )

    # P09 — light decision close; dark is not required to manufacture contrast.
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, PAPER)
    add_rect(slide, 0.58, 0.72, 0.1, 1.08, fill=ORANGE, radius="small")
    add_text(slide, "需要管理层确认", 0.86, 0.92, 5.1, 0.3, size=10, color=ORANGE, bold=True)
    set_slide_title(slide, "今天只做三个决定", 0.86, 1.38, 6.8, 0.7, size=32, color=INK, bold=True)
    decisions = [
        ("01", "将留存恢复设为下季度第一优先级", "同意 / 调整", ORANGE),
        ("02", "授权商务负责人调整渠道折扣边界", "同意 / 调整", BLUE),
        ("03", "按月复核留存、毛利与高质量管道", "同意 / 调整", GREEN),
    ]
    for idx, (num, statement, status, color) in enumerate(decisions):
        y = 2.46 + idx * 1.13
        pale = PALE_ORANGE if color == ORANGE else PALE_BLUE if color == BLUE else PALE_GREEN
        add_rect(slide, 0.72, y - 0.13, 11.45, 0.88, fill=WHITE, line=LINE, radius="medium")
        add_text(slide, num, 1.02, y + 0.12, 0.55, 0.26, size=10, color=color, bold=True)
        add_text(slide, statement, 1.79, y + 0.01, 7.9, 0.54, size=15, color=INK, bold=True)
        add_rect(slide, 9.94, y, 1.85, 0.46, fill=pale, radius="pill")
        add_text(slide, status, 10.05, y + 0.13, 1.63, 0.2, size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.86, 6.13, 5.84, 0.08, fill=ORANGE, radius="small")
    add_text(slide, "下一步：会后 24 小时内锁定责任人与首次复核日期。", 0.86, 6.39, 9.5, 0.3, size=11, color=MUTED)
    add_chrome(slide, 9, "决策与交付")
    add_note(slide, "最后只请管理层做三个决定：第一，是否将留存恢复设为最高优先级；第二，是否授权调整渠道折扣边界；第三，是否按月复核三项核心指标。会后二十四小时内锁定责任人和首次复核日期。")

    return prs


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--no-copy", action="store_true", help="Do not mirror the PPTX into examples/.")
    args = parser.parse_args()
    output = args.output.expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    deck = build_deck()
    deck.save(output)
    normalized_axis_references = normalize_chart_axis_ids(output)
    if not args.no_copy:
        EXAMPLE_DIR.mkdir(parents=True, exist_ok=True)
        mirror = EXAMPLE_DIR / output.name
        if mirror.resolve() != output:
            shutil.copy2(output, mirror)
    print(
        f"Generated {len(deck.slides)} editable slides: {output} "
        f"({normalized_axis_references} signed chart-axis references normalized)"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
